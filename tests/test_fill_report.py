"""Tests for fill outcome tracking (engine.fill_report + fill_template_report)."""

import json
import os

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from engine.fill_report import FillReport, append_fill_history
from engine.pptx_fill import fill_template_report


def _deck_with_boxes(tmp_path, texts):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i, text in enumerate(texts):
        box = slide.shapes.add_textbox(Inches(1), Inches(0.5 + i),
                                       Inches(5), Inches(0.8))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(14)
    path = str(tmp_path / "t.pptx")
    prs.save(path)
    return path


def test_clean_fill_reports_ok(tmp_path):
    tpl = _deck_with_boxes(tmp_path, ["Imps: [I]", "Clicks: [C]"])
    out = str(tmp_path / "o.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        "0": {"assignments": [{"metric": "I", "format": "text",
                               "replace_text": "[I]"}]},
        "1": {"assignments": [{"metric": "C", "format": "text",
                               "replace_text": "[C]"}]},
    }}}}

    _, report = fill_template_report(tpl, out, mapping, {"I": "1", "C": "2"})

    assert report.ok is True
    assert report.filled == 2
    assert report.missing_metrics == []
    assert "Filled 2 text assignment(s)." in report.summary()


def test_missing_metric_and_skip_are_distinguished(tmp_path):
    tpl = _deck_with_boxes(tmp_path, ["A", "B", "C"])
    out = str(tmp_path / "o.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        "0": {"assignments": [{"metric": "Present", "format": "text"}]},
        "1": {"assignments": [{"metric": "Video Starts", "format": "text"}]},
        "2": {"skip": True},
    }}}}

    _, report = fill_template_report(tpl, out, mapping, {"Present": "x"})

    assert report.ok is False
    assert report.filled == 1
    assert report.missing_metrics == ["Video Starts"]
    assert report.skipped == 1
    assert "Video Starts" in report.summary()


def test_unmatched_placeholder_is_the_loud_silent_noop(tmp_path):
    """A replace_text that no longer exists in the template used to vanish
    without a trace — now it's reported."""
    tpl = _deck_with_boxes(tmp_path, ["Impressions: 12,345"])  # already filled
    out = str(tmp_path / "o.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        "0": {"assignments": [{"metric": "I", "format": "text",
                               "replace_text": "[IMP]"}]}}}}}

    _, report = fill_template_report(tpl, out, mapping, {"I": "99"})

    assert report.ok is False
    assert report.filled == 0
    assert report.unmatched_placeholders == [
        {"metric": "I", "placeholder": "[IMP]"}]
    assert "[IMP]" in report.summary()


def test_missing_image_and_out_of_range_counted(tmp_path):
    tpl = _deck_with_boxes(tmp_path, ["LOGO"])
    out = str(tmp_path / "o.pptx")
    mapping = {"slides": {
        "1": {"shape_mappings": {
            "0": {"image_path": "no/such/logo.png"},
            "9": {"assignments": [{"metric": "M", "format": "text"}]}}},
        "7": {"shape_mappings": {"0": {"assignments": [
            {"metric": "M", "format": "text"}]}}},
    }}

    _, report = fill_template_report(tpl, out, mapping, {"M": "x"})

    assert report.missing_images == ["no/such/logo.png"]
    assert report.out_of_range == 2
    assert report.ok is False


def test_failed_query_reported_once_not_as_missing_metric(tmp_path, monkeypatch):
    import engine.query_resolver as qr
    monkeypatch.setattr(qr, "resolve_query",
                        lambda *a, **k: (_ for _ in ()).throw(RuntimeError("x")))
    tpl = _deck_with_boxes(tmp_path, ["Q"])
    out = str(tmp_path / "o.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {"0": {"assignments": [
        {"metric": "Q:imps", "format": "text",
         "query": {"metric": "Impressions"}}]}}}}}

    _, report = fill_template_report(tpl, out, mapping,
                                     {"__client_data__": []})

    assert report.failed_queries == ["Q:imps"]
    assert report.missing_metrics == []  # not double-counted
    assert report.ok is False


def test_fill_template_backcompat_returns_path_only(tmp_path):
    from engine.pptx_fill import fill_template
    tpl = _deck_with_boxes(tmp_path, ["X"])
    out = str(tmp_path / "o.pptx")

    result = fill_template(tpl, out, {"slides": {}}, {})

    assert result == out  # old callers keep working unchanged


def test_history_appends_jsonl(tmp_path, monkeypatch):
    from config import paths
    monkeypatch.setattr(paths, "LOGS_DIR", str(tmp_path / "logs"))
    r1 = FillReport("Deck.pptx", "out1.pptx")
    r1.note_filled()
    r2 = FillReport("Deck.pptx", "out2.pptx")
    r2.note_missing_metric("Reach")

    path = append_fill_history(r1)
    assert append_fill_history(r2) == path

    lines = [json.loads(l) for l in open(path, encoding="utf-8")]
    assert len(lines) == 2
    assert lines[0]["ok"] is True and lines[0]["filled"] == 1
    assert lines[1]["ok"] is False
    assert lines[1]["missing_metrics"] == ["Reach"]
    assert os.path.basename(path) == "fill_history.jsonl"


def test_history_failure_never_raises(monkeypatch):
    from config import paths
    monkeypatch.setattr(paths, "LOGS_DIR", "/proc/definitely/not/writable")
    assert append_fill_history(FillReport()) is None  # swallowed, logged
