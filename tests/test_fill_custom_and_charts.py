"""Classic fill-engine additions (2026-07-17 field feedback):

- Custom text rides IN the saved assignment as a ``{"custom_text": ...}``
  query, so Auto-Fill resolves it in a fresh session — previously the
  text lived only in the mapper session's metric dict and Auto-Fill
  filled nothing.
- Chart assignments store an Apply-as-Chart-Data builder query, and the
  built-in engine injects the resolved categories/series — previously
  chart data only ever reached the deck through the COM live preview.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from engine.pptx_fill import fill_template_report
from engine.query_resolver import resolve_query


def _fresh_session_metrics(client_data=None):
    """What Auto-Fill passes: metrics rebuilt from data — NO session-only
    keys like __custom_<hash>__."""
    return {"__client_data__": client_data or [], "__client_name__": "",
            "__start_date__": "", "__end_date__": ""}


def test_resolve_query_returns_custom_text():
    assert resolve_query({"custom_text": "Hello\nWorld"}, []) == "Hello\nWorld"


def test_custom_text_assignment_survives_to_autofill(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6),
                                   Inches(1))
    box.text_frame.text = "TAGLINE HERE"
    template = str(tmp_path / "t.pptx")
    prs.save(template)

    mapping = {"slides": {"1": {"shape_mappings": {"0": {
        "metric": "__custom_12345__", "format": "text",
        "replace_text": "TAGLINE HERE", "skip": False, "assignments": [],
        "query": {"custom_text": "Powered by Acme"}}}}}}

    out = str(tmp_path / "out.pptx")
    _, report = fill_template_report(template, out, mapping,
                                     _fresh_session_metrics())
    assert report.filled == 1 and report.ok
    deck = Presentation(out)
    text = deck.slides[0].shapes[0].text_frame.text
    assert text == "POWERED BY ACME"   # case matches the placeholder


def test_chart_query_assignment_fills_via_autofill(tmp_path, client_data):
    from lxml import etree
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["OLD-A", "OLD-B"]
    data.add_series("Impressions", (1, 2))
    frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                   Inches(1), Inches(1), Inches(7),
                                   Inches(4), data)
    series = frame.chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0x00, 0x54, 0xA6)
    template = str(tmp_path / "t.pptx")
    prs.save(template)

    chart_query = {"metric": "Impressions", "breakdown": "all",
                   "filter": "all", "agg": "sum", "top_n": "all",
                   "campaigns": [], "sources": ["zip"], "values": [],
                   "output": "chart"}
    mapping = {"slides": {"1": {"shape_mappings": {"0": {
        "metric": "Query: Impressions (sum)", "format": "chart",
        "replace_text": "", "skip": False, "assignments": [],
        "query": chart_query}}}}}

    out = str(tmp_path / "out.pptx")
    _, report = fill_template_report(template, out, mapping,
                                     _fresh_session_metrics(client_data))
    assert report.filled == 1 and report.ok

    deck = Presentation(out)
    chart = deck.slides[0].shapes[0].chart
    assert list(chart.plots[0].categories) == ["33607", "33609"]
    assert list(chart.series[0].values) == [35000.0, 25000.0]
    # Injection replaces cached data only — the brand series color stays
    ser_xml = etree.tostring(chart.series[0]._element).decode()
    assert "0054A6" in ser_xml


def test_chart_query_with_no_matching_rows_reports_loudly(tmp_path):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["A"]
    data.add_series("S", (1,))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1),
                           Inches(1), Inches(6), Inches(4), data)
    template = str(tmp_path / "t.pptx")
    prs.save(template)

    mapping = {"slides": {"1": {"shape_mappings": {"0": {
        "metric": "q", "format": "chart", "replace_text": "", "skip": False,
        "assignments": [],
        "query": {"metric": "Impressions", "agg": "sum", "top_n": "all",
                  "campaigns": [], "sources": ["zip"], "values": [],
                  "output": "chart"}}}}}}

    out = str(tmp_path / "out.pptx")
    _, report = fill_template_report(template, out, mapping,
                                     _fresh_session_metrics([]))
    assert not report.ok
    assert any("matched no data rows" in e for e in report.errors)
