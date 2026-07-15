"""Golden-file characterization tests for engine.pptx_fill.

These lock in the CURRENT observable behavior of fill_template so the
planned MappingModel refactor can be verified against a safety net.
Every test builds a real .pptx with python-pptx, runs the real fill,
re-opens the output, and asserts on what a user would see. No COM, no
display — runs anywhere.
"""

import os

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from engine.pptx_fill import fill_template


# ── Template builders ─────────────────────────────────────────────────────────

def _new_deck(tmp_path, name="template.pptx", slides=1):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(slides):
        prs.slides.add_slide(blank)
    path = str(tmp_path / name)
    prs.save(path)
    return path


def _add_textbox(pptx_path, slide_idx, runs, *, left=1, top=1, width=4, height=1):
    """Append a textbox whose first paragraph contains the given runs.

    ``runs`` is a list of (text, bold) tuples so tests can verify that
    run-level formatting survives the fill.  Returns the shape index the
    fill engine will use (positional, matching current behavior).
    """
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    para = box.text_frame.paragraphs[0]
    for text, bold in runs:
        run = para.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(18)
    shape_idx = len(slide.shapes) - 1
    prs.save(pptx_path)
    return shape_idx


def _shape_text(output_path, slide_idx, shape_idx):
    prs = Presentation(output_path)
    return prs.slides[slide_idx].shapes[shape_idx].text_frame.text


def _mapping(slide_num, shape_idx, **assignment):
    """Build a minimal one-shape mapping in the current on-disk schema."""
    return {"slides": {str(slide_num): {"shape_mappings": {
        str(shape_idx): {"assignments": [assignment]}}}}}


def _tiny_png(path):
    from PIL import Image
    Image.new("RGB", (8, 8), (0, 82, 155)).save(path)
    return str(path)


# ── Full replacement ──────────────────────────────────────────────────────────

def test_full_replacement_sets_text_and_keeps_first_run_formatting(tmp_path):
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("Total Impressions: ", True), ("PLACEHOLDER", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out, _mapping(1, idx, metric="Impressions", format="text"),
                  {"Impressions": "1,234,567"})

    prs = Presentation(out)
    para = prs.slides[0].shapes[idx].text_frame.paragraphs[0]
    # Current contract: whole line replaced, new text lands in run 0,
    # remaining runs are emptied but keep their formatting objects.
    assert para.runs[0].text == "1,234,567"
    assert para.runs[0].font.bold is True
    assert all(r.text == "" for r in para.runs[1:])


def test_full_replacement_targets_first_nonempty_paragraph_only(tmp_path):
    tpl = _new_deck(tmp_path)
    prs = Presentation(tpl)
    box = prs.slides[0].shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = box.text_frame
    tf.paragraphs[0].text = ""            # leading empty paragraph
    tf.add_paragraph().text = "OLD VALUE"  # first content paragraph
    tf.add_paragraph().text = "footnote"   # must be untouched
    shape_idx = len(prs.slides[0].shapes) - 1
    prs.save(tpl)
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out, _mapping(1, shape_idx, metric="M", format="text"),
                  {"M": "NEW"})

    prs = Presentation(out)
    paras = prs.slides[0].shapes[shape_idx].text_frame.paragraphs
    assert paras[1].text == "NEW"
    assert paras[2].text == "footnote"


# ── Partial replacement (replace_text) ────────────────────────────────────────

def test_partial_replacement_within_one_run_preserves_surroundings(tmp_path):
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("Delivered [IMP] impressions", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="Impressions", format="text",
                           replace_text="[IMP]"),
                  {"Impressions": "98,765"})

    assert _shape_text(out, 0, idx) == "Delivered 98,765 impressions"


def test_partial_replacement_spanning_runs_keeps_outer_formatting(tmp_path):
    tpl = _new_deck(tmp_path)
    # Placeholder token split across three runs, bold prefix/suffix around it.
    idx = _add_textbox(tpl, 0, [("Reach: ", True), ("[TO", False),
                                ("TAL]", False), (" households", True)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="Reach", format="text",
                           replace_text="[TOTAL]"),
                  {"Reach": "45,210"})

    prs = Presentation(out)
    para = prs.slides[0].shapes[idx].text_frame.paragraphs[0]
    assert "".join(r.text for r in para.runs) == "Reach: 45,210 households"
    assert para.runs[0].text == "Reach: " and para.runs[0].font.bold is True
    assert para.runs[-1].font.bold is True


def test_partial_replacement_replaces_first_occurrence_only(tmp_path):
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("[X] of [X]", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="M", format="text", replace_text="[X]"),
                  {"M": "7"})

    assert _shape_text(out, 0, idx) == "7 of [X]"


# ── Mapping schema compatibility ──────────────────────────────────────────────

def test_legacy_single_metric_mapping_still_fills(tmp_path):
    """Pre-assignments mappings store metric/format at the shape level."""
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("PLACEHOLDER", False)])
    out = str(tmp_path / "out.pptx")
    legacy = {"slides": {"1": {"shape_mappings": {str(idx): {
        "metric": "Clicks", "format": "text"}}}}}

    fill_template(tpl, out, legacy, {"Clicks": "512"})

    assert _shape_text(out, 0, idx) == "512"


def test_skip_flag_and_missing_metric_leave_shape_untouched(tmp_path):
    tpl = _new_deck(tmp_path)
    idx_a = _add_textbox(tpl, 0, [("KEEP A", False)])
    idx_b = _add_textbox(tpl, 0, [("KEEP B", False)], top=3)
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        str(idx_a): {"skip": True,
                     "assignments": [{"metric": "M", "format": "text"}]},
        str(idx_b): {"assignments": [{"metric": "NotProvided", "format": "text"}]},
    }}}}

    fill_template(tpl, out, mapping, {"M": "SHOULD NOT APPEAR"})

    assert _shape_text(out, 0, idx_a) == "KEEP A"
    assert _shape_text(out, 0, idx_b) == "KEEP B"


def test_out_of_range_slide_and_shape_indices_are_ignored(tmp_path):
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("KEEP", False)])
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {
        "9": {"shape_mappings": {"0": {"assignments": [
            {"metric": "M", "format": "text"}]}}},
        "1": {"shape_mappings": {"99": {"assignments": [
            {"metric": "M", "format": "text"}]}}},
    }}

    fill_template(tpl, out, mapping, {"M": "X"})  # must not raise

    assert _shape_text(out, 0, idx) == "KEEP"


# ── Number / date formatting through the fill path ────────────────────────────

def test_format_details_currency_and_custom_suffix(tmp_path):
    tpl = _new_deck(tmp_path)
    idx_a = _add_textbox(tpl, 0, [("COST", False)])
    idx_b = _add_textbox(tpl, 0, [("RATE", False)], top=3)
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        str(idx_a): {"assignments": [{"metric": "Cost", "format": "number",
                                      "format_details": {"format": "currency",
                                                         "decimals": 2,
                                                         "commas": True}}]},
        str(idx_b): {"assignments": [{"metric": "CompletionRate",
                                      "format": "number",
                                      "format_details": {"format": "custom",
                                                         "decimals": 2,
                                                         "commas": False,
                                                         "prefix": "",
                                                         "suffix": "%"}}]},
    }}}}

    fill_template(tpl, out, mapping,
                  {"Cost": 98765.4321, "CompletionRate": 97.4567})

    assert _shape_text(out, 0, idx_a) == "$98,765.43"
    assert _shape_text(out, 0, idx_b) == "97.46%"


def test_numpy_scalar_values_format_like_native_numbers(tmp_path):
    np = pytest.importorskip("numpy")
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("IMPS", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="Impressions", format="number",
                           format_details={"format": "number", "decimals": 0,
                                           "commas": True}),
                  {"Impressions": np.int64(2500000)})

    assert _shape_text(out, 0, idx) == "2,500,000"


def test_date_range_matches_existing_text_style_and_separator(tmp_path):
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("May 1, 2026 – May 31, 2026", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out, _mapping(1, idx, metric="DateRange", format="date"),
                  {"DateRange": "2026-06-01 - 2026-06-30"})

    # Current behavior: "May 1, 2026" is detected as the ordinal long style,
    # and the en-dash separator from the template is preserved.
    assert _shape_text(out, 0, idx) == "June 1st, 2026 – June 30th, 2026"


def test_partial_replace_uses_target_substring_for_format_context(tmp_path):
    """Regression guard: format detection reads the placeholder being
    replaced, not unrelated placeholders elsewhere in the same line."""
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("[MONTH/YEAR] report for [CLIENT]", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="Client", format="text",
                           replace_text="[CLIENT]"),
                  {"Client": "Acme Motors"})

    # Current behavior: format context comes from the TARGET substring (the
    # fix this test guards), and because "[CLIENT]" is all-caps the engine's
    # case-matching uppercases the replacement. Known quirk: all-caps
    # placeholder tokens force all-caps values.
    assert _shape_text(out, 0, idx) == "[MONTH/YEAR] report for ACME MOTORS"


# ── Image replacement ─────────────────────────────────────────────────────────

def test_image_replacement_preserves_geometry_and_removes_original(tmp_path):
    pytest.importorskip("PIL")
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("LOGO HERE", False)],
                       left=2, top=2, width=3, height=1)
    png = _tiny_png(tmp_path / "logo.png")
    out = str(tmp_path / "out.pptx")
    original = Presentation(tpl).slides[0].shapes[idx]
    geometry = (original.left, original.top, original.width, original.height)

    fill_template(tpl, out,
                  {"slides": {"1": {"shape_mappings": {str(idx): {
                      "image_path": png}}}}},
                  {})

    prs = Presentation(out)
    shapes = list(prs.slides[0].shapes)
    pictures = [s for s in shapes if s.shape_type == 13]  # MSO PICTURE
    assert len(pictures) == 1
    pic = pictures[0]
    assert (pic.left, pic.top, pic.width, pic.height) == geometry
    assert all("LOGO HERE" not in (s.text_frame.text if s.has_text_frame else "")
               for s in shapes)


def test_image_relative_path_resolves_from_app_root(tmp_path, monkeypatch):
    pytest.importorskip("PIL")
    import engine.pptx_fill as pf
    monkeypatch.setattr(pf, "APP_ROOT", str(tmp_path))
    (tmp_path / "templates" / "images").mkdir(parents=True)
    _tiny_png(tmp_path / "templates" / "images" / "logo.png")
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("LOGO", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  {"slides": {"1": {"shape_mappings": {str(idx): {
                      "image_path": os.path.join("templates", "images",
                                                 "logo.png")}}}}},
                  {})

    prs = Presentation(out)
    assert any(s.shape_type == 13 for s in prs.slides[0].shapes)


def test_missing_image_leaves_shape_in_place(tmp_path):
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("KEEP", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  {"slides": {"1": {"shape_mappings": {str(idx): {
                      "image_path": "does/not/exist.png"}}}}},
                  {})  # must not raise

    assert _shape_text(out, 0, idx) == "KEEP"


# ── Query-based assignments ───────────────────────────────────────────────────

def test_query_assignment_resolves_and_caches_value(tmp_path, monkeypatch):
    import engine.query_resolver as qr
    calls = []

    def fake_resolve(query, client_data, client_name, start, end):
        calls.append((query, client_name, start, end))
        return 41250

    monkeypatch.setattr(qr, "resolve_query", fake_resolve)
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("Q", False)])
    out = str(tmp_path / "out.pptx")
    metrics = {"__client_data__": [{"stub": True}],
               "__client_name__": "Acme", "__start_date__": "2026-06-01",
               "__end_date__": "2026-06-30"}
    query = {"metric": "Impressions", "level": "device", "agg": "sum"}

    fill_template(tpl, out,
                  _mapping(1, idx, metric="Q:device_imps", format="number",
                           query=query),
                  metrics)

    assert calls == [(query, "Acme", "2026-06-01", "2026-06-30")]
    assert metrics["Q:device_imps"] == 41250  # cached for reuse
    assert _shape_text(out, 0, idx) == "41,250"


def test_failed_query_resolution_does_not_break_other_shapes(tmp_path, monkeypatch):
    import engine.query_resolver as qr

    def boom(*_a, **_k):
        raise RuntimeError("resolver exploded")

    monkeypatch.setattr(qr, "resolve_query", boom)
    tpl = _new_deck(tmp_path)
    idx_q = _add_textbox(tpl, 0, [("QUERY SHAPE", False)])
    idx_ok = _add_textbox(tpl, 0, [("PLAIN", False)], top=3)
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        str(idx_q): {"assignments": [{"metric": "Q:x", "format": "text",
                                      "query": {"metric": "Impressions"}}]},
        str(idx_ok): {"assignments": [{"metric": "Clicks", "format": "text"}]},
    }}}}

    fill_template(tpl, out, mapping, {"Clicks": "512"})  # must not raise

    assert _shape_text(out, 0, idx_q) == "QUERY SHAPE"
    assert _shape_text(out, 0, idx_ok) == "512"


# ── Multi-slide / multi-assignment ────────────────────────────────────────────

def test_multiple_slides_and_assignments_fill_independently(tmp_path):
    tpl = _new_deck(tmp_path, slides=2)
    idx1 = _add_textbox(tpl, 0, [("Imps: [I], Clicks: [C]", False)])
    idx2 = _add_textbox(tpl, 1, [("Slide two placeholder", False)])
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {
        "1": {"shape_mappings": {str(idx1): {"assignments": [
            {"metric": "I", "format": "text", "replace_text": "[I]"},
            {"metric": "C", "format": "text", "replace_text": "[C]"},
        ]}}},
        "2": {"shape_mappings": {str(idx2): {"assignments": [
            {"metric": "T", "format": "text"}]}}},
    }}

    fill_template(tpl, out, mapping, {"I": "100", "C": "7", "T": "June Recap"})

    assert _shape_text(out, 0, idx1) == "Imps: 100, Clicks: 7"
    assert _shape_text(out, 1, idx2) == "June Recap"  # sentence case: no rule matches, value passes through


def test_text_case_matching_follows_existing_text(tmp_path):
    """Current behavior: full replacement adopts the casing of the text it
    replaces — UPPER stays upper, Title stays title. Single words with
    internal capitals (acronyms like CTR) are exempt."""
    tpl = _new_deck(tmp_path)
    idx_upper = _add_textbox(tpl, 0, [("QUARTERLY TOTAL", False)])
    idx_title = _add_textbox(tpl, 0, [("Quarterly Total", False)], top=3)
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        str(idx_upper): {"assignments": [{"metric": "A", "format": "text"}]},
        str(idx_title): {"assignments": [{"metric": "B", "format": "text"}]},
    }}}}

    fill_template(tpl, out, mapping,
                  {"A": "june recap", "B": "june recap"})

    assert _shape_text(out, 0, idx_upper) == "JUNE RECAP"
    assert _shape_text(out, 0, idx_title) == "June Recap"


# ── Image path fallbacks and failure isolation ────────────────────────────────

def test_image_relative_path_resolves_from_workspace_dir(tmp_path, monkeypatch):
    """Mappings created before the workspace restructure store paths like
    'templates/images/logo.png' relative to workspace/, not the app root."""
    pytest.importorskip("PIL")
    import engine.pptx_fill as pf
    monkeypatch.setattr(pf, "APP_ROOT", str(tmp_path / "nowhere"))
    monkeypatch.setattr(pf, "WORKSPACE_DIR", str(tmp_path / "ws"))
    (tmp_path / "ws" / "templates" / "images").mkdir(parents=True)
    _tiny_png(tmp_path / "ws" / "templates" / "images" / "logo.png")
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("LOGO", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  {"slides": {"1": {"shape_mappings": {str(idx): {
                      "image_path": os.path.join("templates", "images",
                                                 "logo.png")}}}}},
                  {})

    prs = Presentation(out)
    assert any(s.shape_type == 13 for s in prs.slides[0].shapes)


def test_image_path_abs_fallback_when_relative_paths_fail(tmp_path, monkeypatch):
    pytest.importorskip("PIL")
    import engine.pptx_fill as pf
    monkeypatch.setattr(pf, "APP_ROOT", str(tmp_path / "nowhere"))
    monkeypatch.setattr(pf, "WORKSPACE_DIR", str(tmp_path / "nowhere-either"))
    png = _tiny_png(tmp_path / "elsewhere.png")
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("LOGO", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  {"slides": {"1": {"shape_mappings": {str(idx): {
                      "image_path": "gone/logo.png",
                      "image_path_abs": png}}}}},
                  {})

    prs = Presentation(out)
    assert any(s.shape_type == 13 for s in prs.slides[0].shapes)


def test_corrupt_image_error_is_isolated_and_reported(tmp_path):
    """An unreadable image file must not abort the fill: the error lands in
    the report and every other shape still fills."""
    from engine.pptx_fill import fill_template_report

    tpl = _new_deck(tmp_path)
    idx_img = _add_textbox(tpl, 0, [("LOGO", False)])
    idx_txt = _add_textbox(tpl, 0, [("PLAIN", False)], top=3)
    fake_png = tmp_path / "corrupt.png"
    fake_png.write_text("not an image")
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {
        str(idx_img): {"image_path": str(fake_png)},
        str(idx_txt): {"assignments": [{"metric": "M", "format": "text"}]},
    }}}}

    _path, report = fill_template_report(tpl, out, mapping, {"M": "512"})

    assert len(report.errors) == 1
    assert report.ok is False
    assert _shape_text(out, 0, idx_img) == "LOGO"  # original shape kept
    assert _shape_text(out, 0, idx_txt) == "512"


def test_image_assignment_takes_precedence_over_text_assignments(tmp_path):
    """A shape mapped to both an image and text assignments becomes the
    image; the text assignments on that shape are ignored."""
    pytest.importorskip("PIL")
    from engine.pptx_fill import fill_template_report

    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("LOGO", False)])
    png = _tiny_png(tmp_path / "logo.png")
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {str(idx): {
        "image_path": png,
        "assignments": [{"metric": "M", "format": "text"}]}}}}}

    _path, report = fill_template_report(tpl, out, mapping, {"M": "IGNORED"})

    assert report.images_filled == 1
    assert report.filled == 0
    prs = Presentation(out)
    assert all("IGNORED" not in (s.text_frame.text if s.has_text_frame else "")
               for s in prs.slides[0].shapes)


# ── Report-visible fill edge cases ────────────────────────────────────────────

def test_format_details_date_style_formats_range(tmp_path):
    """A date format_details wins over numeric coercion and restyles both
    ends of an ISO range."""
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("REPORT PERIOD", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="__date_range__", format="date",
                           format_details={"format": "date",
                                           "date_style": "us_slash"}),
                  {"__date_range__": "2026-06-01 - 2026-06-30"})

    assert _shape_text(out, 0, idx) == "06/01/2026 – 06/30/2026"


def test_mixed_matched_and_unmatched_placeholders_on_one_shape(tmp_path):
    """One assignment lands, its sibling's placeholder is gone: the fill
    applies what it can and the report names only the unmatched one."""
    from engine.pptx_fill import fill_template_report

    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("Imps: [I], Clicks: [C]", False)])
    out = str(tmp_path / "out.pptx")
    mapping = {"slides": {"1": {"shape_mappings": {str(idx): {"assignments": [
        {"metric": "I", "format": "text", "replace_text": "[I]"},
        {"metric": "X", "format": "text", "replace_text": "[GONE]"},
    ]}}}}}

    _path, report = fill_template_report(tpl, out, mapping,
                                         {"I": "100", "X": "7"})

    assert report.filled == 1
    assert report.unmatched_placeholders == [
        {"metric": "X", "placeholder": "[GONE]"}]
    assert _shape_text(out, 0, idx) == "Imps: 100, Clicks: [C]"


def test_text_assignment_on_shape_without_text_frame_is_silent_noop(tmp_path):
    """Current behavior: a text assignment mapped onto a picture shape does
    nothing and is not reported. Locked here so a future change to surface
    it is deliberate."""
    pytest.importorskip("PIL")
    from engine.pptx_fill import fill_template_report

    tpl = _new_deck(tmp_path)
    png = _tiny_png(tmp_path / "pic.png")
    prs = Presentation(tpl)
    prs.slides[0].shapes.add_picture(png, Inches(1), Inches(1))
    pic_idx = len(prs.slides[0].shapes) - 1
    prs.save(tpl)
    out = str(tmp_path / "out.pptx")

    _path, report = fill_template_report(
        tpl, out, _mapping(1, pic_idx, metric="M", format="text"),
        {"M": "value"})  # must not raise

    assert report.filled == 0
    assert report.ok is True


# ── Multi-line targets and values (Windows batch 4 regressions) ───────────────
# A replace target selected across lines carries '\n' (the UI preview joins
# paragraphs with it) or '\v' (PowerPoint soft line breaks) — characters the
# run text being searched NEVER contains. The replace silently no-oped while
# the report claimed "filled", and multi-line VALUES lost their line breaks
# (a literal '\n' inside <a:t> renders as whitespace, not a break).

def _add_two_paragraph_box(pptx_path, line1, line2):
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1.5))
    tf = box.text_frame
    tf.text = line1
    p2 = tf.add_paragraph()
    p2.text = line2
    shape_idx = len(slide.shapes) - 1
    prs.save(pptx_path)
    return shape_idx


def test_multiline_target_across_paragraphs_fills_and_reports_ok(tmp_path):
    from engine.pptx_fill import fill_template_report
    tpl = _new_deck(tmp_path)
    idx = _add_two_paragraph_box(tpl, "Prepared for: Client Name",
                                 "Presented by: Account Executive")
    out = str(tmp_path / "out.pptx")
    mapping = _mapping(1, idx, metric="AE", format="text",
                       replace_text="Prepared for: Client Name\n"
                                    "Presented by: Account Executive")

    _path, report = fill_template_report(tpl, out, mapping,
                                         {"AE": "Jordan Fox, Spectrum Reach"})

    assert report.ok is True and report.filled == 1
    text = _shape_text(out, 0, idx)
    assert "Jordan Fox, Spectrum Reach" in text
    assert "Account Executive" not in text  # later target lines cleared


def test_soft_linebreak_target_fills(tmp_path):
    """One paragraph broken with <a:br/>: scan/preview text carries '\\v',
    which the runs never contain."""
    from engine.pptx_fill import fill_template_report
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "Client Name\vAccount Executive"  # -> run <a:br/> run
    tpl = str(tmp_path / "t.pptx")
    prs.save(tpl)
    out = str(tmp_path / "out.pptx")
    mapping = _mapping(1, 0, metric="AE", format="text",
                       replace_text="Client Name\vAccount Executive")

    _path, report = fill_template_report(tpl, out, mapping, {"AE": "Jordan"})

    assert report.ok is True
    text = _shape_text(out, 0, 0)
    assert "Jordan" in text and "Account Executive" not in text


def test_failed_partial_replace_is_reported_not_filled(tmp_path):
    from engine.pptx_fill import fill_template_report
    tpl = _new_deck(tmp_path)
    idx = _add_textbox(tpl, 0, [("Static text", False)])
    out = str(tmp_path / "out.pptx")
    mapping = _mapping(1, idx, metric="M", format="text",
                       replace_text="Nope line 1\nNope line 2")

    _path, report = fill_template_report(tpl, out, mapping, {"M": "42"})

    assert report.ok is False
    assert report.filled == 0
    assert report.unmatched_placeholders  # loudly reported, never silent


def test_multiline_value_becomes_real_line_breaks(tmp_path):
    """A value containing '\\n' must produce <a:br/> elements — a literal
    newline inside <a:t> renders as whitespace in PowerPoint."""
    tpl = _new_deck(tmp_path)
    # Mixed-case context so the all-caps case-forcing quirk stays out of
    # this test's way — the line-break behavior is what's under test.
    idx = _add_textbox(tpl, 0, [("Contact: ", True), ("[Ph]", False)])
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out,
                  _mapping(1, idx, metric="C", format="text",
                           replace_text="[Ph]"),
                  {"C": "Jordan Fox\n555-0100"})

    prs = Presentation(out)
    para = prs.slides[0].shapes[idx].text_frame.paragraphs[0]
    from pptx.oxml.ns import qn
    assert len(para._p.findall(qn("a:br"))) == 1     # real line break
    assert all("\n" not in r.text for r in para.runs)  # no literal newlines
    assert "Jordan Fox" in para.text and "555-0100" in para.text
    assert para.runs[0].font.bold is True  # formatting still intact


def test_full_assignment_to_empty_box_writes_value(tmp_path):
    """Assigning a full-text value to an empty box used to no-op silently."""
    from engine.pptx_fill import fill_template_report
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tpl = str(tmp_path / "t.pptx")
    prs.save(tpl)
    out = str(tmp_path / "out.pptx")

    _path, report = fill_template_report(tpl, out,
                                         _mapping(1, 0, metric="M",
                                                  format="text"),
                                         {"M": "FILLED"})

    assert report.filled == 1
    assert _shape_text(out, 0, 0) == "FILLED"
