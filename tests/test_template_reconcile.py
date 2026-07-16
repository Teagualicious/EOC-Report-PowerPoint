"""Template-first Phase D tests: re-ingest reconciliation (review work
carries forward across template revisions; only deltas surface) and the
workflow-level AI tools (ingest_template_store / build_template_report).
"""

import os

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from engine.template_ir import (classify_template, ingest_template,
                                load_template_ir, new_slot_mapping,
                                reingest_template, save_slot_mapping,
                                save_template_ir, validate_slot_mapping)
from engine.template_ir.classify import (rename_slot, set_classification,
                                         set_excluded, shape_slots)
from engine.template_ir.reconcile import PREV_JSON_NAME, reconcile


def _build_deck_v1(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    Inches(10), Inches(0.6))
    banner.name = "Banner"
    header = slide.shapes.add_textbox(Inches(0.3), Inches(0.1),
                                      Inches(8), Inches(0.5))
    header.name = "Header"
    header.text_frame.text = "CLIENT NAME | MONTH 1ST, 2026"
    title = slide.shapes.add_textbox(Inches(0.3), Inches(1),
                                     Inches(8), Inches(0.6))
    title.name = "Title"
    title.text_frame.text = "Extending Your Brand's Reach"
    kpi = slide.shapes.add_textbox(Inches(3), Inches(2), Inches(4),
                                   Inches(1.2))
    kpi.name = "KPI"
    run = kpi.text_frame.paragraphs[0].add_run()
    run.text = "X,XXX"
    run.font.size = Pt(48)
    label = kpi.text_frame.add_paragraph().add_run()
    label.text = "IMPRESSIONS"
    prs.save(path)


def _edit_deck_to_v2(v1_path, v2_path):
    """The client's revision: title copy changed, header removed, a new
    completion-rate callout added. KPI and banner untouched."""
    prs = Presentation(v1_path)
    slide = prs.slides[0]
    for shape in list(slide.shapes):
        if shape.name == "Header":
            shape._element.getparent().remove(shape._element)
        elif shape.name == "Title":
            shape.text_frame.text = "A Fresh New Tagline"
    new_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(3),
                                       Inches(1))
    new_box.name = "Completion"
    r = new_box.text_frame.paragraphs[0].add_run(); r.text = "XX%"
    lab = new_box.text_frame.add_paragraph().add_run()
    lab.text = "COMPLETION RATE"
    prs.save(v2_path)


@pytest.fixture
def reviewed_store(tmp_path):
    """v1 ingested + classified + human-reviewed: slot renamed, title
    forced dynamic, banner excluded, mapping saved."""
    v1 = str(tmp_path / "deck_v1.pptx")
    _build_deck_v1(v1)
    store = ingest_template(v1, template_id="evolve",
                            store_dir=str(tmp_path / "store"))
    ir = classify_template(load_template_ir(store))
    rename_slot(ir, "impressions", "total_impressions")
    title = next(s for slide in ir.slides for s in slide.shapes
                 if s.name == "Title")
    set_classification(ir, title.shape_id, "dynamic")   # override: dynamic
    banner = next(s for slide in ir.slides for s in slide.shapes
                  if s.name == "Banner")
    set_excluded(ir, banner.shape_id, True)
    save_template_ir(ir, store)

    mapping = new_slot_mapping("evolve")
    mapping["slots"] = {
        "total_impressions": {"query": "__total_Impressions__",
                              "format": "number"},
        "client_name": {"query": "__client_name__", "format": "text"},
    }
    save_slot_mapping(mapping, store)
    return tmp_path, v1, store


def _deltas_by_kind(deltas):
    by_kind = {}
    for d in deltas:
        by_kind.setdefault(d["kind"], []).append(d)
    return by_kind


def test_reingest_carries_review_work_and_reports_deltas(reviewed_store):
    tmp_path, v1, store = reviewed_store
    v2 = str(tmp_path / "deck_v2.pptx")
    _edit_deck_to_v2(v1, v2)

    _, deltas = reingest_template(v2, store)
    ir = load_template_ir(store)
    by_name = {s.name: s for slide in ir.slides for s in slide.shapes}
    kinds = _deltas_by_kind(deltas)

    # Review work carried: rename, exclusion, and the KPI slot's anchor
    assert "total_impressions" in ir.slot_registry
    spec = ir.slot_registry["total_impressions"]
    assert spec["shape_uid"] == by_name["KPI"].shape_uid
    assert spec["placeholder_text"] == "X,XXX"
    assert by_name["Banner"].excluded
    assert by_name["Title"].classification == "dynamic"   # override kept

    # Deltas: removed header (its mapped slot flagged), changed title
    # (its slot's placeholder is gone), added shape with a new slot
    assert any(d["shape"] for d in kinds["shape_removed"])
    dropped = {d["slot"]: d["detail"] for d in kinds["slot_dropped"]}
    assert "client_name" in dropped
    assert "WAS mapped" in dropped["client_name"]
    assert any("placeholder text no longer appears" in reason
               for reason in dropped.values())            # the title slot
    assert any(d["shape"] == by_name["Completion"].shape_id
               for d in kinds["shape_added"])
    assert "completion_rate" in ir.slot_registry
    assert any(d["slot"] == "completion_rate" for d in kinds["slot_new"])
    assert any(d["shape"] == by_name["Title"].shape_id
               for d in kinds["shape_text_changed"])

    # Safety net + surviving mapping entry still validates cleanly
    assert os.path.isfile(os.path.join(store, PREV_JSON_NAME))
    warnings = validate_slot_mapping(ir, {"slots": {
        "total_impressions": {"query": "__total_Impressions__"}}})
    assert not any("total_impressions" in w and "not a slot" in w
                   for w in warnings)


def test_reingest_of_identical_deck_is_a_quiet_no_op(reviewed_store):
    tmp_path, v1, store = reviewed_store
    before = load_template_ir(store)
    _, deltas = reingest_template(v1, store)
    after = load_template_ir(store)
    assert deltas == []
    assert list(after.slot_registry) == list(before.slot_registry)
    assert [s.classification for slide in after.slides for s in slide.shapes] \
        == [s.classification for slide in before.slides for s in slide.shapes]


def test_unreviewed_store_gets_fresh_suggestions_only(tmp_path):
    """A Phase A store (no classifications, no slots) has no review work
    to carry — re-ingest must not smear its all-static defaults over the
    fresh suggestions."""
    v1 = str(tmp_path / "deck.pptx")
    _build_deck_v1(v1)
    store = ingest_template(v1, store_dir=str(tmp_path / "store"))
    old_ir = load_template_ir(store)                 # never classified
    new_ir = classify_template(load_template_ir(store))
    deltas = reconcile(old_ir, new_ir)
    assert [d["kind"] for d in deltas] == ["no_prior_review"]
    assert "impressions" in new_ir.slot_registry     # suggestions intact
    kpi = next(s for slide in new_ir.slides for s in slide.shapes
               if s.name == "KPI")
    assert kpi.classification == "dynamic"


def test_duplicate_names_never_match(tmp_path):
    """Two shapes sharing a name are ambiguous — reconciliation must not
    guess between them (the Phase 4 rule)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(2):
        box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(3),
                                       Inches(0.5))
        box.name = "Twin"
        box.text_frame.text = f"X,XXX{i}"
    deck = str(tmp_path / "twins.pptx")
    prs.save(deck)
    store = ingest_template(deck, store_dir=str(tmp_path / "store"))
    ir = classify_template(load_template_ir(store))
    old_shapes = [s for slide_ir in ir.slides for s in slide_ir.shapes]

    new_ir = classify_template(load_template_ir(store))
    # Strip uids so only the (duplicate) names could match
    for slide_ir in new_ir.slides:
        for s in slide_ir.shapes:
            s.shape_uid = None
    for s in old_shapes:
        s.shape_uid = None
    deltas = reconcile(ir, new_ir)
    kinds = {d["kind"] for d in deltas}
    assert "shape_added" in kinds and "shape_removed" in kinds


# ── Workflow-level AI tools ───────────────────────────────────────────────────

@pytest.fixture
def isolated_store(tmp_path, monkeypatch):
    from engine.template_ir import ingest
    store_root = str(tmp_path / "template_store")
    monkeypatch.setattr(ingest, "TEMPLATE_STORE_DIR", store_root)
    return store_root


def test_workflow_ingest_and_reingest(tmp_path, isolated_store):
    from engine import workflow
    v1 = str(tmp_path / "Client Deck.pptx")
    _build_deck_v1(v1)

    result = workflow.ingest_template_store(v1)
    assert result["template_id"] == "client_deck"
    assert result["deltas"] == []
    assert "impressions" in result["slots"]
    assert any(s["classification"] == "dynamic" for s in result["shapes"])

    stores = workflow.list_template_stores()
    assert [s["template_id"] for s in stores] == ["client_deck"]
    assert stores[0]["has_mapping"] is False

    v2 = str(tmp_path / "deck_v2.pptx")
    _edit_deck_to_v2(v1, v2)
    result = workflow.ingest_template_store(v2, template_id="client_deck")
    kinds = {d["kind"] for d in result["deltas"]}
    assert {"shape_removed", "shape_added", "slot_new"} <= kinds


def test_workflow_build_template_report(tmp_path, isolated_store,
                                        client_data):
    from engine import workflow
    v1 = str(tmp_path / "Client Deck.pptx")
    _build_deck_v1(v1)
    result = workflow.ingest_template_store(v1)
    store = result["template_dir"]

    mapping = new_slot_mapping("client_deck")
    mapping["slots"] = {
        "impressions": {"query": "__total_Impressions__", "format": "number"},
        "client_name": {"query": "__client_name__", "format": "text"},
    }
    save_slot_mapping(mapping, store)

    out = str(tmp_path / "monthly.pptx")
    built = workflow.build_template_report(client_data, "client_deck", out,
                                           client_name="Acme Appliance Co")
    assert built["path"] == out and os.path.isfile(out)
    assert built["report"]["slots_filled"] == 2
    deck = Presentation(out)
    texts = " ".join(s.text_frame.text for s in deck.slides[0].shapes
                     if s.has_text_frame)
    assert "148,450" in texts and "ACME APPLIANCE CO" in texts

    with pytest.raises(FileNotFoundError):
        workflow.build_template_report(client_data, "no_such_store", out)


def test_cli_ingest_template_and_stores(tmp_path, isolated_store, capsys):
    import json
    import cli
    v1 = str(tmp_path / "Client Deck.pptx")
    _build_deck_v1(v1)

    code = cli.main(["ingest-template", "--pptx", v1])
    payload = json.loads(capsys.readouterr().out)
    assert code == 0 and payload["ok"]
    assert "impressions" in payload["data"]["slots"]

    code = cli.main(["template-stores"])
    payload = json.loads(capsys.readouterr().out)
    assert code == 0 and payload["ok"]
    assert payload["data"][0]["template_id"] == "client_deck"
