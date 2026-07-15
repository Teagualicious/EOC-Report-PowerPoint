"""Unit tests for mapper.mapping_model.MappingModel — the single owner of
template-mapping state (mapper roadmap Phase 3).

The model must preserve the exact on-disk schema and the assignment
semantics the pre-model UI implemented inline, because engine.pptx_fill
and saved user mappings both depend on them.
"""

import json

import pytest

from mapper.mapping_model import (APPENDED, CREATED, NEEDS_CONFIRM, REPLACED,
                                  UPDATED, MappingModel)


@pytest.fixture
def model():
    return MappingModel()


# ── assign_metric semantics ───────────────────────────────────────────────────

def test_first_assignment_uses_legacy_shape_level_schema(model):
    outcome = model.assign_metric(1, 3, "Impressions", fmt="number")

    assert outcome == CREATED
    smap = model.data["slides"]["1"]["shape_mappings"]["3"]
    # Contract with engine.pptx_fill: first assignment lives at shape level
    # with an empty assignments list (the legacy single-assignment form).
    assert smap["metric"] == "Impressions"
    assert smap["format"] == "number"
    assert smap["skip"] is False
    assert smap["assignments"] == []


def test_partial_assignment_appends_to_existing(model):
    model.assign_metric(1, 3, "Impressions")
    outcome = model.assign_metric(1, 3, "Clicks", replace_text="[C]")

    assert outcome == APPENDED
    asgns = model.assignments(1, 3)
    assert [a["metric"] for a in asgns] == ["Impressions", "Clicks"]
    assert asgns[1]["replace_text"] == "[C]"
    # The legacy shape-level entry was normalized into the list form
    assert model.shape_map(1, 3).get("metric") is None


def test_same_metric_updates_in_place_and_keeps_siblings(model):
    model.assign_metric(1, 3, "Impressions", replace_text="[I]")
    model.assign_metric(1, 3, "Clicks", replace_text="[C]")

    details = {"format": "currency", "decimals": 0, "commas": True,
               "prefix": "", "suffix": ""}
    outcome = model.assign_metric(1, 3, "Clicks", fmt="currency",
                                  format_details=details)

    assert outcome == UPDATED
    asgns = model.assignments(1, 3)
    assert len(asgns) == 2
    clicks = next(a for a in asgns if a["metric"] == "Clicks")
    assert clicks["format"] == "currency"
    assert clicks["format_details"] == details
    # No new replace target given — the existing one is kept
    assert clicks["replace_text"] == "[C]"


def test_full_replace_requires_confirmation(model):
    model.assign_metric(1, 3, "Impressions", replace_text="[I]")

    outcome = model.assign_metric(1, 3, "Clicks")  # full-text, no confirm
    assert outcome == NEEDS_CONFIRM
    # Nothing changed
    assert [a["metric"] for a in model.assignments(1, 3)] == ["Impressions"]

    outcome = model.assign_metric(1, 3, "Clicks", confirm_replace=True)
    assert outcome == REPLACED
    smap = model.shape_map(1, 3)
    assert smap["metric"] == "Clicks"
    assert smap["assignments"] == []


def test_legacy_mapping_normalized_on_append_preserving_fields():
    legacy = {"slides": {"1": {"shape_mappings": {"0": {
        "metric": "Impressions", "format": "number", "replace_text": "[I]",
        "query": {"metric": "Impressions"}, "last_rendered": "1,204,556",
    }}}}}
    model = MappingModel(legacy)

    model.assign_metric(1, 0, "Clicks", replace_text="[C]")

    asgns = model.assignments(1, 0)
    first = asgns[0]
    assert first["metric"] == "Impressions"
    assert first["replace_text"] == "[I]"
    assert first["query"] == {"metric": "Impressions"}
    assert first["last_rendered"] == "1,204,556"
    assert asgns[1]["metric"] == "Clicks"


# ── Image / skip / clear ──────────────────────────────────────────────────────

def test_assign_image_schema(model):
    model.assign_metric(1, 3, "Impressions")
    model.assign_image(1, 3, "templates/images/logo.png", "/abs/logo.png")

    smap = model.shape_map(1, 3)
    assert smap == {"metric": "", "image_path": "templates/images/logo.png",
                    "image_path_abs": "/abs/logo.png", "skip": False,
                    "assignments": []}


def test_set_skip_preserves_assignments(model):
    """Owner decision 2026-07-15 (mapper roadmap Phase 5): Skip is a flag
    on top of the mapping — toggling it must never discard assignments.
    (Replaces the pre-model discard behavior this test previously locked.)"""
    model.assign_metric(1, 3, "Impressions", replace_text="[I]")
    model.set_skip(1, 3, True)

    assert model.shape_map(1, 3)["skip"] is True
    assert [a["metric"] for a in model.assignments(1, 3)] == ["Impressions"]

    model.set_skip(1, 3, False)
    assert model.shape_map(1, 3)["skip"] is False
    assert [a["metric"] for a in model.assignments(1, 3)] == ["Impressions"]

    model.set_skip(2, 9, True)  # never-mapped shape — creates a bare flag
    assert model.shape_map(2, 9) == {"skip": True}


def test_clear_shape_removes_mapping_and_tolerates_unknown(model):
    model.assign_metric(1, 3, "Impressions")
    model.clear_shape(1, 3)
    assert model.shape_map(1, 3) == {}

    model.clear_shape(2, 9)  # never mapped — must not raise


# ── Per-metric format preferences ─────────────────────────────────────────────

def test_set_metric_format_propagates_to_all_assignment_forms(model):
    details = {"format": "currency", "decimals": 2, "commas": True,
               "prefix": "", "suffix": ""}
    model.assign_metric(1, 0, "Spend")                       # legacy shape-level
    model.assign_metric(2, 5, "Spend", replace_text="[S]")   # shape-level too
    model.assign_metric(2, 5, "Clicks", replace_text="[C]")  # list form now
    model.assign_metric(3, 1, "Clicks")                      # untouched metric

    touched = model.set_metric_format("Spend", "currency", details)

    assert touched == 2
    assert model.metric_formats["Spend"] == "currency"
    assert model.metric_format_details["Spend"] == details
    assert model.shape_map(1, 0)["format_details"] == details
    assert model.shape_map(1, 0)["format"] == "currency"
    spend = next(a for a in model.assignments(2, 5) if a["metric"] == "Spend")
    assert spend["format_details"] == details
    clicks = next(a for a in model.assignments(2, 5) if a["metric"] == "Clicks")
    assert "format_details" not in clicks


def test_set_metric_format_without_details_only_stores_preference(model):
    model.assign_metric(1, 0, "Spend")
    touched = model.set_metric_format("Spend", "number")

    assert touched == 0
    assert model.metric_formats["Spend"] == "number"
    assert "format_details" not in model.shape_map(1, 0)


# ── note_rendered bookkeeping ─────────────────────────────────────────────────

def test_note_rendered_writes_both_forms_silently(model):
    events = []
    model.subscribe(lambda event, **kw: events.append(event))
    model.assign_metric(1, 0, "Impressions")                   # shape-level
    model.assign_metric(2, 5, "Impressions", replace_text="[I]")
    model.assign_metric(2, 5, "Clicks", replace_text="[C]")
    events.clear()

    model.note_rendered(1, 0, "Impressions", "1,204,556")
    model.note_rendered(2, 5, "Impressions", "1,204,556")

    assert model.shape_map(1, 0)["last_rendered"] == "1,204,556"
    imp = next(a for a in model.assignments(2, 5)
               if a["metric"] == "Impressions")
    assert imp["last_rendered"] == "1,204,556"
    clk = next(a for a in model.assignments(2, 5) if a["metric"] == "Clicks")
    assert "last_rendered" not in clk
    assert events == []  # bookkeeping never notifies


# ── Observers ─────────────────────────────────────────────────────────────────

def test_observers_receive_each_mutation(model):
    events = []
    model.subscribe(lambda event, slide_num=None, shape_id=None, metric=None:
                    events.append((event, slide_num, shape_id, metric)))

    model.assign_metric(1, 3, "Impressions")
    model.assign_image(1, 4, "rel.png", "/abs.png")
    model.set_skip(1, 5, True)
    model.clear_shape(1, 3)
    model.set_metric_format("Impressions", "number",
                            {"format": "number", "decimals": 0,
                             "commas": True, "prefix": "", "suffix": ""})

    assert events == [
        ("assign", 1, 3, "Impressions"),
        ("image", 1, 4, None),
        ("skip", 1, 5, None),
        ("clear", 1, 3, None),
        ("format", None, None, "Impressions"),
    ]


def test_failing_observer_does_not_break_mutations(model):
    def boom(event, **kw):
        raise RuntimeError("observer exploded")
    model.subscribe(boom)

    outcome = model.assign_metric(1, 3, "Impressions")  # must not raise

    assert outcome == CREATED
    assert model.shape_map(1, 3)["metric"] == "Impressions"


# ── Scan-identity stamping (mapper roadmap Phase 4) ───────────────────────────

SCAN = [{"slide_num": 1, "shapes": [
    {"shape_id": 0, "shape_uid": 7, "name": "TextBox 3"},
    {"shape_id": 1, "shape_uid": None, "name": "No-Uid Shape"},
]}]


def test_mutations_stamp_scanned_identity(model):
    model.set_scan_identity(SCAN)

    model.assign_metric(1, 0, "Impressions", replace_text="[I]")
    smap = model.shape_map(1, 0)
    assert smap["shape_uid"] == 7
    assert smap["shape_name"] == "TextBox 3"

    model.assign_image(1, 0, "rel.png", "/abs.png")
    smap = model.shape_map(1, 0)
    assert smap["shape_uid"] == 7 and smap["shape_name"] == "TextBox 3"

    model.set_skip(1, 0, True)  # skip targets a specific shape too
    smap = model.shape_map(1, 0)
    assert smap["skip"] is True
    assert smap["shape_uid"] == 7 and smap["shape_name"] == "TextBox 3"
    assert smap["image_path"] == "rel.png"  # skip no longer discards


def test_shape_without_uid_stays_unstamped(model):
    model.set_scan_identity(SCAN)
    model.assign_metric(1, 1, "Clicks")
    smap = model.shape_map(1, 1)
    assert "shape_uid" not in smap and "shape_name" not in smap


def test_no_scan_identity_keeps_schema_byte_identical(model):
    """Callers that never attach a scan (headless flows, every pre-Phase-4
    code path) must see the exact legacy schema."""
    model.assign_metric(1, 0, "Impressions")
    model.set_skip(1, 5, True)
    assert "shape_uid" not in model.shape_map(1, 0)
    assert model.shape_map(1, 5) == {"skip": True}


def test_legacy_loaded_entry_upgrades_on_touch():
    """Entries from an old saved mapping gain identity lazily — the first
    time the user touches them — and never by merely loading/saving."""
    legacy = {"slides": {"1": {"shape_mappings": {"0": {
        "metric": "Impressions", "format": "number", "replace_text": "[I]",
    }}}}}
    model = MappingModel(legacy)
    model.set_scan_identity(SCAN)
    assert "shape_uid" not in model.shape_map(1, 0)  # loading didn't stamp

    model.assign_metric(1, 0, "Impressions", fmt="number")  # touch
    assert model.shape_map(1, 0)["shape_uid"] == 7


# ── Persistence and fill-engine integration ───────────────────────────────────

def test_to_dict_is_json_serializable_and_loads_back(model):
    model.assign_metric(1, 0, "Impressions", fmt="number")
    model.assign_metric(1, 0, "Clicks", replace_text="[C]")
    model.assign_image(1, 2, "templates/images/logo.png")
    model.set_skip(1, 5, True)

    reloaded = MappingModel(json.loads(json.dumps(model.to_dict())))

    assert reloaded.to_dict() == model.to_dict()
    assert [a["metric"] for a in reloaded.assignments(1, 0)] == \
           ["Impressions", "Clicks"]


def test_model_built_mapping_drives_fill_template(tmp_path):
    """End-to-end: a mapping assembled through the model fills a real deck
    exactly like a hand-written mapping dict."""
    from pptx import Presentation
    from pptx.util import Inches
    from engine.pptx_fill import fill_template

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Imps: [I]"
    tpl = str(tmp_path / "tpl.pptx")
    prs.save(tpl)

    model = MappingModel()
    model.assign_metric(1, 0, "Impressions", fmt="number", replace_text="[I]")
    out = str(tmp_path / "out.pptx")

    fill_template(tpl, out, model.to_dict(), {"Impressions": 1204556})

    filled = Presentation(out)
    assert filled.slides[0].shapes[0].text_frame.text == "Imps: 1,204,556"
