"""Tests for the AI-facing layer: engine/workflow.py + app/cli.py.

The workflow service must produce the same numbers as the interactive
Quick Run flow — it IS the same engine code, exercised end to end here
with synthetic exports and an isolated workspace.
"""

import json
import os

import pytest

from conftest import PROJECT_ROOT

ARCHITECT_CONFIG = {"sheets": [
    {"sheet_name": "Data by Devices", "columns": [
        {"column": "Order", "role": "campaign_id", "selected": True},
        {"column": "Device", "role": "level: device", "selected": True},
        {"column": "Impressions", "role": "metric", "selected": True},
        {"column": "100% Complete", "role": "metric", "selected": True},
    ]},
    {"sheet_name": "Data by Geo", "columns": [
        {"column": "Order", "role": "campaign_id", "selected": True},
        {"column": "Zip Code", "role": "level: zip", "selected": True},
        {"column": "Impressions", "role": "metric", "selected": True},
    ]},
]}


@pytest.fixture
def isolated_workspace(tmp_path, monkeypatch):
    """Point settings/mappings at a temp workspace with one platform saved."""
    from config import paths, settings
    ws = tmp_path / "workspace"
    (ws / "mappings").mkdir(parents=True)
    (ws / "templates").mkdir()
    monkeypatch.setattr(paths, "WORKSPACE_DIR", str(ws))
    monkeypatch.setattr(paths, "SETTINGS_PATH", str(ws / "settings.json"))
    monkeypatch.setattr(paths, "MAPPINGS_DIR", str(ws / "mappings"))
    monkeypatch.setattr(paths, "TEMPLATES_DIR", str(ws / "templates"))
    monkeypatch.setattr(settings, "SETTINGS_PATH", str(ws / "settings.json"))
    monkeypatch.setattr(settings, "MAPPINGS_DIR", str(ws / "mappings"))
    from engine import workflow
    monkeypatch.setattr(workflow, "TEMPLATES_DIR", str(ws / "templates"))
    settings.save_settings({"platforms": {"Architect": True}})
    settings.save_platform_config("Architect", ARCHITECT_CONFIG)
    return ws


class TestWorkflowService:
    def test_parse_campaigns_kpis_match_pipeline(self, isolated_workspace,
                                                 sample_xlsx):
        from engine import workflow

        parsed, failures = workflow.parse_files([(sample_xlsx, "Architect")])
        assert failures == []
        campaigns = workflow.list_campaigns(parsed)
        assert {"name": "Campaign A", "platform": "Architect"} in campaigns

        data = workflow.client_dataset(parsed, "Acme Appliance Co",
                                       start_date="2026-05-01",
                                       end_date="2026-05-31")
        kpis = workflow.kpis_for(data)
        # Best source per campaign: A zip 40,100 beats devices 40,000;
        # B devices 40,000 beat zip 39,900 -> total 80,100
        assert kpis["totals"]["Impressions"] == 80100
        assert kpis["campaign_details"]["Campaign A"]["Impressions"] == 40100

    def test_parse_failure_reported_not_raised(self, isolated_workspace,
                                               sample_csv):
        from engine import workflow
        parsed, failures = workflow.parse_files(
            [("/nowhere/nothing.csv", "Architect"),
             (sample_csv, "NoSuchPlatform")])
        assert parsed == []
        assert ("nothing.csv", "File not found") in failures
        assert any("NoSuchPlatform" in reason for _n, reason in failures)

    def test_campaign_filter_limits_dataset(self, isolated_workspace,
                                            sample_xlsx):
        from engine import workflow
        parsed, _ = workflow.parse_files([(sample_xlsx, "Architect")])
        data = workflow.client_dataset(parsed, "FT", campaigns=["Campaign A"])
        kpis = workflow.kpis_for(data)
        assert "Campaign B" not in kpis["campaign_details"]

    def test_export_workbook_writes_search_dashboard(self, isolated_workspace,
                                                     sample_xlsx, tmp_path):
        from engine import workflow
        from openpyxl import load_workbook

        parsed, _ = workflow.parse_files([(sample_xlsx, "Architect")])
        data = workflow.client_dataset(parsed, "Acme Appliance Co")
        result = workflow.export_workbook(data, "Acme Appliance Co",
                                          str(tmp_path / "out"),
                                          inject_vba=False)
        assert os.path.exists(result["path"])
        # Same folder/file naming as the interactive export (safe name
        # helpers keep spaces)
        assert "Acme Appliance Co_unified_report" in result["path"]
        wb = load_workbook(result["path"])
        assert {"Search", "Unified Data"} <= set(wb.sheetnames)
        assert result["rows"] > 0

    def test_fill_deck_fills_mapped_template(self, isolated_workspace,
                                             sample_xlsx, tmp_path):
        from engine import workflow
        from engine.pptx_mapper import save_template_mapping
        from pptx import Presentation
        from pptx.util import Inches

        # A one-shape template mapped to the client name, saved in the
        # isolated workspace exactly as the mapper would
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1),
                                       Inches(4), Inches(1))
        # Mixed case on purpose: all-caps template text triggers the
        # locked case-forcing quirk (would render "ACME APPLIANCE CO")
        box.text_frame.text = "Placeholder"
        tpl_name = "unit_template.pptx"
        prs.save(str(isolated_workspace / "templates" / tpl_name))
        save_template_mapping(tpl_name, {"slides": {"1": {"shape_mappings": {
            "0": {"metric": "__client_name__", "format": "text",
                  "skip": False, "assignments": []}}}}})

        parsed, _ = workflow.parse_files([(sample_xlsx, "Architect")])
        data = workflow.client_dataset(parsed, "Acme Appliance Co",
                                       start_date="2026-05-01",
                                       end_date="2026-05-31")
        out = str(tmp_path / "deck.pptx")
        result = workflow.fill_deck(data, tpl_name, out, "Acme Appliance Co")

        assert result["report"]["ok"] is True
        assert result["report"]["filled"] == 1
        filled = Presentation(result["path"])
        assert filled.slides[0].shapes[0].text_frame.text == "Acme Appliance Co"

    def test_query_metric_matches_resolver(self, isolated_workspace,
                                           sample_xlsx):
        from engine import workflow
        parsed, _ = workflow.parse_files([(sample_xlsx, "Architect")])
        data = workflow.client_dataset(parsed, "FT")
        result = workflow.query_metric(data, "Impressions", breakdown="zip")
        assert result["value"] == 80000  # 22,000 + 18,100 + 39,900


class TestCli:
    def _run(self, argv, capsys):
        import cli
        code = cli.main(argv)
        out = capsys.readouterr().out
        return code, json.loads(out)

    def test_platforms_lists_configured(self, isolated_workspace, capsys):
        code, payload = self._run(["platforms"], capsys)
        assert code == 0 and payload["ok"] is True
        assert payload["data"][0]["name"] == "Architect"
        assert payload["data"][0]["sheets"] == 2

    def test_kpis_end_to_end(self, isolated_workspace, sample_xlsx, capsys):
        code, payload = self._run(
            ["kpis", "--file", f"{sample_xlsx}=Architect",
             "--client", "Acme Appliance Co",
             "--start", "2026-05-01", "--end", "2026-05-31"], capsys)
        assert code == 0 and payload["ok"] is True
        # JSON numbers, not stringified numpy scalars
        assert payload["data"]["totals"]["Impressions"] == 80100
        assert isinstance(payload["data"]["totals"]["Impressions"], int)

    def test_error_is_json_with_nonzero_exit(self, isolated_workspace,
                                             capsys):
        code, payload = self._run(
            ["kpis", "--file", "/nowhere.csv=Architect", "--client", "X"],
            capsys)
        assert code == 1
        assert payload["ok"] is False
        assert "No files could be parsed" in payload["error"]

    def test_scan_reports_structure(self, isolated_workspace, sample_xlsx,
                                    capsys):
        code, payload = self._run(["scan", "--file", sample_xlsx], capsys)
        assert code == 0
        names = [s["name"] for s in payload["data"]]
        assert "Data by Devices" in names


def test_cli_and_mcp_share_the_workflow_engine():
    """The MCP server and CLI must stay thin shells over engine.workflow —
    neither may grow its own parsing/aggregation logic."""
    import inspect
    import cli
    from engine import workflow as wf
    src = inspect.getsource(cli)
    assert "from engine import workflow" in src
    for fn in ("parse_files", "client_dataset", "kpis_for",
               "export_workbook", "fill_deck", "query_metric",
               "list_campaigns", "list_platforms", "list_templates",
               "scan_export"):
        assert hasattr(wf, fn)
    mcp_src = open(os.path.join(PROJECT_ROOT, "app", "mcp_server.py"),
                   encoding="utf-8").read()
    assert "from engine import workflow" in mcp_src
