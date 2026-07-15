"""Template-first pipeline Phase A tests (engine.template_ir).

The core invariant from the review doc: ingest → build → re-ingest must
yield an equivalent template, and every build emits a report — silent
partial output is never acceptable. All pure python-pptx; no COM.
"""

import json
import os

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from engine.template_ir import (build_from_template, ingest_template,
                                load_template_ir)
from engine.template_ir.ingest import strip_embed_ids


@pytest.fixture
def branded_deck(tmp_path):
    """A synthetic 'client template': banner rectangle, two-run styled
    textbox, table, logo image — the shape mix Phase A must round-trip."""
    pytest.importorskip("PIL")
    from PIL import Image

    logo = tmp_path / "logo.png"
    Image.new("RGB", (16, 16), (0, 84, 166)).save(logo)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    prs.slide_width, Inches(0.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0x00, 0x54, 0xA6)
    banner.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1.5))
    para = box.text_frame.paragraphs[0]
    r1 = para.add_run(); r1.text = "X,XXX"
    r1.font.size = Pt(48); r1.font.bold = True
    r1.font.color.rgb = RGBColor(0x2E, 0x8B, 0x57)
    r2 = para.add_run(); r2.text = " IMPRESSIONS"
    r2.font.size = Pt(18); r2.font.bold = True

    slide.shapes.add_table(2, 3, Inches(1), Inches(4), Inches(6), Inches(1.5))
    slide.shapes.add_picture(str(logo), Inches(11), Inches(0.1),
                             Inches(1.2), Inches(0.4))

    path = str(tmp_path / "Branded Template.pptx")
    prs.save(path)
    return path


def _store(tmp_path):
    return str(tmp_path / "store")


# ── Ingest ────────────────────────────────────────────────────────────────────

def test_ingest_writes_schema_and_assets(tmp_path, branded_deck):
    template_dir = ingest_template(branded_deck, store_dir=_store(tmp_path))

    assert os.path.isfile(os.path.join(template_dir, "template.json"))
    ir = load_template_ir(template_dir)
    assert ir.schema_version == "1.0"
    assert ir.slide_width_emu == int(Inches(13.333))

    (slide,) = ir.slides
    types = sorted(s.shape_type for s in slide.shapes)
    assert types == ["image", "shape", "table", "text_box"]

    # Persistent identity from day 1 (Phase 4 lesson)
    assert all(isinstance(s.shape_uid, int) for s in slide.shapes)
    assert len({s.shape_uid for s in slide.shapes}) == len(slide.shapes)

    # Verbatim XML stored for every shape; image extracted + recorded
    assert all(s.element_xml.strip().startswith("<") for s in slide.shapes)
    (pic,) = [s for s in slide.shapes if s.shape_type == "image"]
    assert len(pic.image_rels) == 1
    (asset_rel,) = pic.image_rels.values()
    assert os.path.isfile(os.path.join(template_dir, asset_rel))

    # Text snapshot + EMU geometry captured for preview/classify phases
    (text_shape,) = [s for s in slide.shapes if s.shape_type == "text_box"]
    assert "IMPRESSIONS" in text_shape.text
    assert text_shape.geometry["left"] == 914400  # Inches(1)


def test_chart_shapes_are_flagged_and_skipped(tmp_path):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(4), Inches(0.6))
    box.text_frame.text = "KEEP ME"
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("Impressions", (1, 2))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1),
                           Inches(1), Inches(6), Inches(4), data)
    deck = str(tmp_path / "chart_deck.pptx")
    prs.save(deck)

    template_dir = ingest_template(deck, store_dir=_store(tmp_path))
    ir = load_template_ir(template_dir)
    (chart,) = [s for s in ir.slides[0].shapes if s.shape_type == "chart"]
    assert chart.unsupported  # Phase C, loudly

    out = str(tmp_path / "rebuilt.pptx")
    _, report = build_from_template(template_dir, out)
    assert report["shapes_copied"] == 1          # the textbox
    assert [s for s, _ in report["skipped"]] == [chart.shape_id]

    rebuilt = Presentation(out)
    texts = [s.text_frame.text for s in rebuilt.slides[0].shapes
             if s.has_text_frame]
    assert "KEEP ME" in texts


# ── Build: verbatim fidelity ──────────────────────────────────────────────────

def test_build_reproduces_deck_verbatim(tmp_path, branded_deck):
    template_dir = ingest_template(branded_deck, store_dir=_store(tmp_path))
    out = str(tmp_path / "rebuilt.pptx")

    _, report = build_from_template(template_dir, out)

    assert report["skipped"] == []
    assert report["shapes_copied"] == 4
    assert report["images_relinked"] == 1

    original = Presentation(branded_deck)
    rebuilt = Presentation(out)
    assert rebuilt.slide_width == original.slide_width
    assert len(rebuilt.slides[0].shapes) == len(original.slides[0].shapes)

    # Formatting survives because the XML is the SAME: two styled runs
    orig_box = [s for s in original.slides[0].shapes
                if s.has_text_frame and "IMPRESSIONS" in s.text_frame.text][0]
    new_box = [s for s in rebuilt.slides[0].shapes
               if s.has_text_frame and "IMPRESSIONS" in s.text_frame.text][0]
    for orig_run, new_run in zip(orig_box.text_frame.paragraphs[0].runs,
                                 new_box.text_frame.paragraphs[0].runs):
        assert new_run.text == orig_run.text
        assert new_run.font.size == orig_run.font.size
        assert new_run.font.bold == orig_run.font.bold

    # Image bytes are identical after relinking
    orig_pic = [s for s in original.slides[0].shapes
                if s.shape_type == 13][0]  # MSO_SHAPE_TYPE.PICTURE
    new_pic = [s for s in rebuilt.slides[0].shapes if s.shape_type == 13][0]
    assert new_pic.image.blob == orig_pic.image.blob
    assert new_pic.left == orig_pic.left and new_pic.width == orig_pic.width


def test_round_trip_is_idempotent(tmp_path, branded_deck):
    """ingest(build(ingest(deck))) ≡ ingest(deck) — the Phase A invariant."""
    first_dir = ingest_template(branded_deck, template_id="rt",
                                store_dir=str(tmp_path / "s1"))
    rebuilt = str(tmp_path / "rebuilt.pptx")
    build_from_template(first_dir, rebuilt)
    second_dir = ingest_template(rebuilt, template_id="rt",
                                 store_dir=str(tmp_path / "s2"))

    ir1 = load_template_ir(first_dir)
    ir2 = load_template_ir(second_dir)
    assert len(ir2.slides) == len(ir1.slides)
    shapes1 = sorted(ir1.slides[0].shapes, key=lambda s: s.z_order)
    shapes2 = sorted(ir2.slides[0].shapes, key=lambda s: s.z_order)
    assert len(shapes2) == len(shapes1)
    for s1, s2 in zip(shapes1, shapes2):
        assert s2.shape_type == s1.shape_type
        assert s2.geometry == s1.geometry
        assert s2.text == s1.text
        # The XML itself must survive byte-for-byte apart from packaging
        # relationship ids, which legitimately change across rebuilds
        assert strip_embed_ids(s2.element_xml) == strip_embed_ids(s1.element_xml)


def test_schema_json_round_trips_through_disk(tmp_path, branded_deck):
    template_dir = ingest_template(branded_deck, store_dir=_store(tmp_path))
    with open(os.path.join(template_dir, "template.json"), encoding="utf-8") as f:
        raw = json.load(f)
    ir = load_template_ir(template_dir)
    assert ir.to_dict() == raw
