"""Template-first pipeline Phase B tests: classification suggestions,
slot registry, slot mappings, and dynamic text build.

All pure python-pptx/lxml — no COM, no Tk. The build-side invariant:
only slot run TEXT changes; every other byte stays the verbatim Phase A
copy, and every unfilled slot is reported (never silent).
"""

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from engine.template_ir import (build_from_template, build_mapped_report,
                                classify_template, ingest_template,
                                load_slot_mapping, load_template_ir,
                                new_slot_mapping, resolve_slot_values,
                                save_slot_mapping, save_template_ir,
                                validate_slot_mapping)
from engine.template_ir.classify import (looks_like_placeholder, rename_slot,
                                         set_classification, set_excluded,
                                         shape_slots)
from engine.template_ir.schema import TemplateIR


@pytest.fixture
def placeholder_deck(tmp_path):
    """A synthetic client template with the Phase B shape mix: header
    field, KPI callout (value+label runs), grouped KPI, static title,
    branding bar."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    prs.slide_width, Inches(0.6))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0x00, 0x54, 0xA6)

    header = slide.shapes.add_textbox(Inches(0.3), Inches(0.05),
                                      Inches(9), Inches(0.5))
    header.text_frame.text = "CLIENT NAME | MONTH 1ST, 2026"

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.8),
                                     Inches(9), Inches(0.6))
    title.text_frame.text = "Extending Your Brand's Reach"

    kpi = slide.shapes.add_textbox(Inches(4), Inches(1.8), Inches(5),
                                   Inches(1.3))
    p0 = kpi.text_frame.paragraphs[0]
    value_run = p0.add_run()
    value_run.text = "X,XXX"
    value_run.font.size = Pt(48)
    value_run.font.bold = True
    value_run.font.color.rgb = RGBColor(0x2E, 0x8B, 0x57)
    p1 = kpi.text_frame.add_paragraph()
    label_run = p1.add_run()
    label_run.text = "IMPRESSIONS"
    label_run.font.size = Pt(18)

    group = slide.shapes.add_group_shape()
    gbox = group.shapes.add_textbox(Inches(0.5), Inches(4), Inches(3),
                                    Inches(1))
    gp = gbox.text_frame.paragraphs[0]
    gr = gp.add_run(); gr.text = "XX%"
    gl = gbox.text_frame.add_paragraph().add_run()
    gl.text = "COMPLETION RATE"

    path = str(tmp_path / "Placeholder Template.pptx")
    prs.save(path)
    return path


@pytest.fixture
def classified_dir(tmp_path, placeholder_deck):
    """Ingested + classified template store for the deck above."""
    template_dir = ingest_template(placeholder_deck,
                                   store_dir=str(tmp_path / "store"))
    ir = classify_template(load_template_ir(template_dir))
    save_template_ir(ir, template_dir)
    return template_dir


# ── Placeholder heuristics ────────────────────────────────────────────────────

@pytest.mark.parametrize("text,expected_type", [
    ("X,XXX", "number"),
    ("$X,XXX", "number"),
    ("XX%", "number"),
    ("X.X", "number"),
    ("N/A", "text"),
    ("{{total_impressions}}", "text"),
    ("[Client]", "text"),
    ("CLIENT NAME | MONTH 1ST, 2026", "text"),
    ("May 1st, 2026", "date"),
])
def test_placeholder_like_text_is_detected(text, expected_type):
    reason, slot_type = looks_like_placeholder(text)
    assert reason is not None
    assert slot_type == expected_type


@pytest.mark.parametrize("text", [
    "IMPRESSIONS",              # all-caps label — the static half of a pair
    "Extending Your Brand's Reach",
    "XL",                       # letters after X are not a mask
    "", "   ",
])
def test_ordinary_text_is_not_a_placeholder(text):
    assert looks_like_placeholder(text) == (None, None)


# ── Auto-classification + slot registry ──────────────────────────────────────

def test_classify_suggests_dynamic_and_names_slots_from_labels(classified_dir):
    ir = load_template_ir(classified_dir)
    by_text = {s.text: s for s in ir.slides[0].shapes}

    kpi = by_text["X,XXX\nIMPRESSIONS"]
    assert kpi.classification == "dynamic"
    assert kpi.slot_name == "impressions"        # named from the paired label

    header = by_text["CLIENT NAME | MONTH 1ST, 2026"]
    assert header.classification == "dynamic"
    assert header.slot_name == "client_name"

    assert by_text["Extending Your Brand's Reach"].classification == "static"
    banner = next(s for s in ir.slides[0].shapes if s.shape_type == "shape")
    assert banner.classification == "static"        # branding rectangle

    slots = ir.slot_registry
    assert set(slots) == {"impressions", "client_name", "completion_rate"}
    assert slots["impressions"]["type"] == "number"
    assert slots["impressions"]["placeholder_text"] == "X,XXX"
    assert slots["impressions"]["shape_uid"] == kpi.shape_uid
    assert slots["completion_rate"]["type"] == "number"  # found inside group


def test_classification_survives_disk_round_trip(classified_dir):
    ir = load_template_ir(classified_dir)
    assert ir.slot_registry            # persisted by save_template_ir
    assert any(s.classification == "dynamic" for s in ir.slides[0].shapes)


def test_review_overrides_toggle_slots(classified_dir):
    ir = load_template_ir(classified_dir)
    title = next(s for s in ir.slides[0].shapes
                 if s.text == "Extending Your Brand's Reach")
    set_classification(ir, title.shape_id, "dynamic")
    assert title.classification == "dynamic"
    # No placeholder-looking runs — the first non-empty run becomes the slot
    (slot,) = shape_slots(ir, title.shape_id)
    assert ir.slot_registry[slot]["placeholder_text"] == \
        "Extending Your Brand's Reach"

    set_classification(ir, title.shape_id, "static")
    assert shape_slots(ir, title.shape_id) == []
    assert title.slot_name is None


def test_rename_slot_updates_registry_and_shape(classified_dir):
    ir = load_template_ir(classified_dir)
    order_before = list(ir.slot_registry)
    new = rename_slot(ir, "impressions", "Total Impressions!")
    assert new == "total_impressions"
    assert list(ir.slot_registry) == [
        "total_impressions" if n == "impressions" else n
        for n in order_before]                       # order preserved
    kpi = next(s for s in ir.slides[0].shapes if "X,XXX" in s.text)
    assert kpi.slot_name == "total_impressions"

    with pytest.raises(ValueError):
        rename_slot(ir, "client_name", "total_impressions")
    with pytest.raises(ValueError):
        rename_slot(ir, "client_name", "!!!")


def test_v10_template_json_still_loads():
    """Phase A stores (schema 1.0, no classification fields) load as-is."""
    data = {"schema_version": "1.0", "template_id": "old",
            "slide_width_emu": 1, "slide_height_emu": 1,
            "slides": [{"slide_index": 0, "shapes": [{
                "shape_id": "s0_u2", "shape_uid": 2, "name": "Box",
                "shape_type": "text_box", "z_order": 0, "geometry": {},
                "element_xml": "<sp/>", "text": "hello"}]}]}
    ir = TemplateIR.from_dict(data)
    assert ir.slot_registry == {}
    shape = ir.slides[0].shapes[0]
    assert shape.classification == "static" and not shape.excluded


# ── Slot mapping: persistence, validation, resolution ────────────────────────

def _mapping_for(template_dir):
    ir = load_template_ir(template_dir)
    mapping = new_slot_mapping(ir.template_id)
    mapping["slots"] = {
        "impressions": {"query": "__total_Impressions__", "format": "number"},
        "client_name": {"query": "__client_name__", "format": "text"},
    }
    return ir, mapping


def test_mapping_round_trips_through_disk(classified_dir):
    _, mapping = _mapping_for(classified_dir)
    save_slot_mapping(mapping, classified_dir)
    assert load_slot_mapping(classified_dir) == mapping


def test_validation_flags_unknown_unmapped_and_type_mismatch(classified_dir):
    ir, mapping = _mapping_for(classified_dir)
    mapping["slots"]["bogus"] = {"query": "__client_name__"}
    warnings = validate_slot_mapping(ir, mapping)
    assert any("bogus" in w and "not a slot" in w for w in warnings)
    assert any("completion_rate" in w and "no data source" in w
               for w in warnings)

    # A text source mapped onto the number slot is caught at map time
    mapping["slots"]["impressions"]["query"] = "__client_name__"
    warnings = validate_slot_mapping(ir, mapping)
    assert any("impressions" in w and "expects a number" in w
               for w in warnings)


def test_resolve_formats_values_against_placeholder_context(classified_dir,
                                                            client_data):
    ir, mapping = _mapping_for(classified_dir)
    values, issues = resolve_slot_values(ir, mapping, client_data,
                                         client_name="Acme Appliance Co")
    assert issues == []
    assert values["impressions"] == "148,450"          # number formatting
    assert values["client_name"] == "ACME APPLIANCE CO"  # matches deck case


def test_resolve_reports_slots_without_a_source(classified_dir, client_data):
    ir, mapping = _mapping_for(classified_dir)
    mapping["slots"]["completion_rate"] = {"query": None}
    _, issues = resolve_slot_values(ir, mapping, client_data)
    assert ("completion_rate", "no data source mapped") in issues


def test_builder_query_resolves_through_pivot(classified_dir, client_data):
    """Owner decision 2026-07-16: slots reuse the FULL query schema —
    a saved Advanced Query Builder query re-resolves via the same pivot
    the builder displayed."""
    ir, mapping = _mapping_for(classified_dir)
    mapping["slots"]["impressions"]["query"] = {
        "metric": "Impressions", "agg": "sum", "top_n": "all",
        "campaigns": ["Campaign A"], "sources": ["device"], "values": []}
    values, issues = resolve_slot_values(ir, mapping, client_data)
    assert issues == []
    assert values["impressions"] == "50,000"   # device rows of Campaign A


# ── Dynamic text build ────────────────────────────────────────────────────────

def test_build_fills_slots_and_preserves_run_formatting(classified_dir,
                                                        client_data, tmp_path):
    _, mapping = _mapping_for(classified_dir)
    mapping["slots"]["completion_rate"] = {"query": "__total_Impressions__",
                                           "format": "number"}
    save_slot_mapping(mapping, classified_dir)

    out = str(tmp_path / "monthly.pptx")
    _, report = build_mapped_report(classified_dir, out, client_data,
                                    client_name="Acme Appliance Co")
    assert report["slots_filled"] == 3
    assert report["slots_unfilled"] == []
    assert report["skipped"] == []

    deck = Presentation(out)
    texts = {}
    def collect(shapes):
        for s in shapes:
            if s.shape_type == 13 or not hasattr(s, "has_text_frame"):
                continue
            if s.shape_type == 6:            # MSO group
                collect(s.shapes); continue
            if s.has_text_frame:
                texts[s.text_frame.text] = s
    collect(deck.slides[0].shapes)

    kpi = next(s for t, s in texts.items() if "148,450" in t)
    value_run = kpi.text_frame.paragraphs[0].runs[0]
    assert value_run.text == "148,450"
    assert value_run.font.size == Pt(48) and value_run.font.bold  # verbatim
    assert kpi.text_frame.paragraphs[1].runs[0].text == "IMPRESSIONS"

    assert any("ACME APPLIANCE CO" in t for t in texts)   # header filled
    assert any("COMPLETION RATE" in t for t in texts)     # group filled
    assert any("Extending Your Brand's Reach" in t for t in texts)  # static


def test_resolution_issue_for_unknown_slot_still_reported(classified_dir,
                                                          client_data,
                                                          tmp_path):
    """A mapped slot the template no longer knows (no query, not in the
    registry) must surface in the build report, not vanish."""
    _, mapping = _mapping_for(classified_dir)
    mapping["slots"]["ghost"] = {"query": None}
    save_slot_mapping(mapping, classified_dir)
    _, report = build_mapped_report(classified_dir, str(tmp_path / "o.pptx"),
                                    client_data,
                                    client_name="Acme Appliance Co")
    assert ("ghost", "no data source mapped") in report["slots_unfilled"]
    assert ("completion_rate", "no value provided") in report["slots_unfilled"]


def test_unfilled_slots_are_reported_never_silent(classified_dir, tmp_path):
    out = str(tmp_path / "out.pptx")
    _, report = build_from_template(
        classified_dir, out,
        slot_values={"impressions": "1,234", "bogus_slot": "x"})
    assert report["slots_filled"] == 1
    unfilled = dict(report["slots_unfilled"])
    assert unfilled["bogus_slot"] == "not a slot in this template"
    assert unfilled["client_name"] == "no value provided"
    assert unfilled["completion_rate"] == "no value provided"


def test_stale_placeholder_is_reported(classified_dir, tmp_path):
    ir = load_template_ir(classified_dir)
    ir.slot_registry["impressions"]["placeholder_text"] = "Y,YYY"
    save_template_ir(ir, classified_dir)
    _, report = build_from_template(classified_dir, str(tmp_path / "o.pptx"),
                                    slot_values={"impressions": "9"})
    assert ("impressions", "placeholder text not found in shape") \
        in report["slots_unfilled"]


def test_excluded_shape_is_skipped_and_its_slot_reported(classified_dir,
                                                         tmp_path):
    ir = load_template_ir(classified_dir)
    kpi = next(s for s in ir.slides[0].shapes if "X,XXX" in s.text)
    set_excluded(ir, kpi.shape_id, True)
    save_template_ir(ir, classified_dir)

    out = str(tmp_path / "out.pptx")
    _, report = build_from_template(classified_dir, out,
                                    slot_values={"impressions": "1,234"})
    assert (kpi.shape_id, "excluded in template review") in report["skipped"]
    assert ("impressions", "shape excluded in template review") \
        in report["slots_unfilled"]
    deck = Presentation(out)
    assert all("IMPRESSIONS" not in (s.text_frame.text if s.has_text_frame
                                     else "")
               for s in deck.slides[0].shapes if s.shape_type != 6)


# ── Phase C: chart and table slots ────────────────────────────────────────────

@pytest.fixture
def chart_table_dir(tmp_path):
    """Ingested + classified store for a deck with a styled chart and a
    branded table."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    data = CategoryChartData()
    data.categories = ["OLD-A", "OLD-B", "OLD-C"]
    data.add_series("Impressions", (10, 20, 30))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1),
                                Inches(1), Inches(6), Inches(3.5), data)
    gf.name = "Zip Chart"
    series = gf.chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0x00, 0x54, 0xA6)

    table_frame = slide.shapes.add_table(2, 2, Inches(8), Inches(1),
                                         Inches(4), Inches(2))
    table_frame.name = "Delivery Table"
    table = table_frame.table
    table.cell(0, 0).text = "ZIP"
    table.cell(0, 1).text = "IMPRESSIONS"
    table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True

    deck = str(tmp_path / "Chart Table Template.pptx")
    prs.save(deck)
    template_dir = ingest_template(deck, store_dir=str(tmp_path / "ct_store"))
    ir = classify_template(load_template_ir(template_dir))
    save_template_ir(ir, template_dir)
    return template_dir


def _zip_query(output):
    """A saved Advanced Query Builder query over the conftest client_data's
    zip breakdown, applied as the given output kind."""
    return {"metric": "Impressions", "breakdown": "all", "filter": "all",
            "agg": "sum", "top_n": "all", "campaigns": [],
            "sources": ["zip"], "values": [], "output": output}


def test_classify_creates_chart_and_table_slots(chart_table_dir):
    ir = load_template_ir(chart_table_dir)
    slots = ir.slot_registry
    assert slots["zip_chart"]["type"] == "chart_data"
    assert slots["delivery_table"]["type"] == "table_data"
    for shape in ir.slides[0].shapes:
        assert shape.classification == "dynamic"
        assert shape.unsupported is None
    (chart,) = [s for s in ir.slides[0].shapes if s.shape_type == "chart"]
    assert chart.chart_part["xml"] and chart.chart_part["workbook"]


def test_builder_query_resolves_to_payloads(client_data):
    from engine.query_resolver import resolve_query_payload
    table = resolve_query_payload(_zip_query("table"), client_data, "table")
    assert table["headers"] == ["level_value", "Campaign A", "Total"]
    assert table["rows"] == [["33607", "35,000", "35,000"],
                             ["33609", "25,000", "25,000"]]

    chart = resolve_query_payload(_zip_query("chart"), client_data, "chart")
    assert chart["categories"] == ["33607", "33609"]
    assert chart["series"] == [{"name": "Campaign A",
                                "values": [35000.0, 25000.0]}]  # no Total


def test_validation_requires_matching_output_kind(chart_table_dir):
    ir = load_template_ir(chart_table_dir)
    mapping = new_slot_mapping(ir.template_id)
    mapping["slots"] = {"zip_chart": {"query": "__client_name__"},
                        "delivery_table": {"query": _zip_query("value")}}
    warnings = validate_slot_mapping(ir, mapping)
    assert any("zip_chart" in w and "Chart Data" in w for w in warnings)
    assert any("delivery_table" in w and "Table" in w for w in warnings)


def test_build_injects_chart_data_preserving_series_color(chart_table_dir,
                                                          client_data,
                                                          tmp_path):
    from lxml import etree
    mapping = new_slot_mapping("ct")
    mapping["slots"] = {"zip_chart": {"query": _zip_query("chart")}}
    save_slot_mapping(mapping, chart_table_dir)

    out = str(tmp_path / "built.pptx")
    _, report = build_mapped_report(chart_table_dir, out, client_data,
                                    client_name="Acme Appliance Co")
    assert report["charts_cloned"] == 1
    filled = dict.fromkeys([s for s, _ in report["slots_unfilled"]])
    assert "zip_chart" not in filled

    deck = Presentation(out)
    (chart_shape,) = [s for s in deck.slides[0].shapes if s.has_chart]
    chart = chart_shape.chart
    assert list(chart.plots[0].categories) == ["33607", "33609"]
    assert list(chart.series[0].values) == [35000.0, 25000.0]
    # Series-level styling survives injection (replace_data leaves it alone)
    ser_xml = etree.tostring(chart.series[0]._element).decode()
    assert "0054A6" in ser_xml
    # The embedded workbook was replaced too — "Edit Data" stays truthful
    assert chart.part.chart_workbook.xlsx_part is not None


def test_build_fills_table_keeping_branded_header(chart_table_dir,
                                                  client_data, tmp_path):
    mapping = new_slot_mapping("ct")
    mapping["slots"] = {"delivery_table": {"query": _zip_query("table")}}
    save_slot_mapping(mapping, chart_table_dir)

    out = str(tmp_path / "built.pptx")
    _, report = build_mapped_report(chart_table_dir, out, client_data,
                                    client_name="Acme Appliance Co")
    assert ("delivery_table",
            "table has 2 column(s); 1 data column(s) dropped"
            ) not in report["slots_unfilled"]
    assert any("delivery_table" in n and "dropped" in n
               for n in report["notes"])   # Total column didn't fit — loud

    deck = Presentation(out)
    (table_shape,) = [s for s in deck.slides[0].shapes if s.has_table]
    table = table_shape.table
    assert len(table.rows) == 3            # header + 2 data (one row cloned)
    assert table.cell(0, 0).text == "ZIP"  # branded header kept
    assert table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold
    assert table.cell(1, 0).text == "33607"
    assert table.cell(1, 1).text == "35,000"
    assert table.cell(2, 0).text == "33609"


def test_chart_round_trip_without_data_keeps_original_values(chart_table_dir,
                                                             tmp_path):
    """A plain rebuild (no slot values) must still produce a working chart
    showing the template's original cached data."""
    out = str(tmp_path / "plain.pptx")
    _, report = build_from_template(chart_table_dir, out)
    assert report["charts_cloned"] == 1 and report["skipped"] == []
    deck = Presentation(out)
    (chart_shape,) = [s for s in deck.slides[0].shapes if s.has_chart]
    assert list(chart_shape.chart.plots[0].categories) == \
        ["OLD-A", "OLD-B", "OLD-C"]


def test_multiline_value_becomes_real_line_breaks(classified_dir, tmp_path):
    out = str(tmp_path / "out.pptx")
    build_from_template(classified_dir, out,
                        slot_values={"client_name": "Acme Motors\nMay 2026"})
    deck = Presentation(out)
    header = next(s for s in deck.slides[0].shapes
                  if s.has_text_frame and "Acme Motors" in s.text_frame.text)
    # A literal \n in <a:t> renders as whitespace in PowerPoint — the fill
    # must produce a real <a:br/> (reuses pptx_fill._explode_newlines)
    assert "\n" not in header.text_frame.paragraphs[0].runs[0].text
    assert "Acme Motors" in header.text_frame.text
    assert "May 2026" in header.text_frame.text
