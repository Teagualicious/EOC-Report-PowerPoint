"""Regression tests for Windows DPI sizing, template preview fallbacks,
and the background-work helper."""

import time
from pathlib import Path


class _FakeRoot:
    """Minimal root for run_in_background: collects after() callbacks so the
    test drives the poll loop synchronously."""

    def __init__(self):
        self.pending = []

    def after(self, _ms, fn):
        self.pending.append(fn)

    def _drive(self, done, timeout=5.0):
        deadline = time.time() + timeout
        while not done and time.time() < deadline:
            if self.pending:
                self.pending.pop(0)()
            else:
                time.sleep(0.01)


def test_run_in_background_without_progress_keeps_old_contract():
    from ui import utils

    root = _FakeRoot()
    done = []
    utils.run_in_background(root, lambda: 41 + 1, done.append,
                            lambda exc: done.append(exc))
    root._drive(done)

    assert done == [42]


def test_run_in_background_delivers_progress_in_order_before_result():
    from ui import utils

    root = _FakeRoot()
    seen, done = [], []

    def work(progress):
        progress("writing workbook…")
        progress("computing KPIs…")
        return "result"

    utils.run_in_background(root, work, done.append,
                            lambda exc: done.append(exc),
                            on_progress=seen.append)
    root._drive(done)

    assert done == ["result"]
    assert seen == ["writing workbook…", "computing KPIs…"]


def test_fit_window_scales_geometry_for_windows_dpi(monkeypatch):
    from ui import utils

    class FakeTk:
        def call(self, *args):
            assert args == ("tk", "scaling")
            return 2.0  # 150% relative to Tk's 96-DPI baseline of 4/3

    class FakeWindow:
        tk = FakeTk()

        def __init__(self):
            self.value = None

        def geometry(self, value):
            self.value = value

        def winfo_screenwidth(self):
            return 1920

        def winfo_screenheight(self):
            return 1080

    win = FakeWindow()
    monkeypatch.setattr(utils.sys, "platform", "win32")
    monkeypatch.setattr(utils, "_work_area", lambda _win: (0, 0, 1920, 1040))

    utils.fit_window(win, 560, 520)

    assert win.value.startswith("840x780+")
    _, position = win.value.split("+", 1)
    x, y = (int(v) for v in position.split("+"))
    assert x >= 0
    assert y >= 0


def test_fit_window_clamps_scaled_dialog_to_work_area(monkeypatch):
    from ui import utils

    class FakeTk:
        def call(self, *args):
            return 2.4

    class FakeWindow:
        tk = FakeTk()

        def __init__(self):
            self.value = None

        def geometry(self, value):
            self.value = value

        def winfo_screenwidth(self):
            return 1280

        def winfo_screenheight(self):
            return 720

    win = FakeWindow()
    monkeypatch.setattr(utils.sys, "platform", "win32")
    monkeypatch.setattr(utils, "_work_area", lambda _win: (0, 0, 1280, 680))

    utils.fit_window(win, 940, 700)

    size = win.value.split("+", 1)[0]
    width, height = (int(v) for v in size.split("x"))
    assert width < 1280
    assert height < 680


def test_template_preview_uses_text_when_thumbnail_unavailable(tmp_path, monkeypatch):
    from pptx import Presentation
    from engine import pptx_thumbs

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 2_000_000, 500_000)
    box.text_frame.text = "Monthly Performance Summary"
    prs.save(path)

    monkeypatch.setattr(pptx_thumbs, "get_template_thumbnail",
                        lambda *_args, **_kwargs: None)

    preview = pptx_thumbs.get_template_preview(str(path))

    assert preview["kind"] == "text"
    assert "Monthly Performance Summary" in preview["value"]


def test_template_preview_returns_cached_image_path(tmp_path, monkeypatch):
    from engine import pptx_thumbs

    template = tmp_path / "sample.pptx"
    template.write_bytes(b"not needed for mocked thumbnail")
    png = Path(tmp_path / "sample.png")
    png.write_bytes(b"png")
    monkeypatch.setattr(pptx_thumbs, "get_template_thumbnail",
                        lambda *_args, **_kwargs: str(png))

    preview = pptx_thumbs.get_template_preview(str(template))

    assert preview == {"kind": "image", "value": str(png)}
