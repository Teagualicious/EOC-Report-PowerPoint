"""Characterization tests for mapping persistence and template scanning.

Covers engine.pptx_mapper: save/load round-trips, corrupt-file handling,
template listing, and the python-pptx scan structure the mapper UI is
built on.  Locks the on-disk schema before the MappingModel refactor.
"""

import json
import os

import pytest
from pptx import Presentation
from pptx.util import Inches

from config import paths


@pytest.fixture
def store(tmp_path, monkeypatch):
    """Point mapping/template storage at a temp dir (same pattern as
    test_template_bundle)."""
    tdir = tmp_path / "templates"
    mdir = tmp_path / "mappings"
    tdir.mkdir()
    mdir.mkdir()
    monkeypatch.setattr(paths, "TEMPLATES_DIR", str(tdir))
    monkeypatch.setattr(paths, "MAPPINGS_DIR", str(mdir))
    import engine.pptx_mapper as pm
    monkeypatch.setattr(pm, "TEMPLATES_DIR", str(tdir))
    monkeypatch.setattr(pm, "MAPPINGS_DIR", str(mdir))
    return tmp_path


SAMPLE_MAPPING = {"slides": {"1": {"shape_mappings": {
    "0": {"assignments": [{"metric": "Impressions", "format": "number",
                           "format_details": {"format": "number",
                                              "decimals": 0, "commas": True},
                           "replace_text": "[IMP]"}]},
    "3": {"skip": True},
    "5": {"image_path": "templates/images/logo.png"},
}}}}


# ── Persistence round-trip ────────────────────────────────────────────────────

def test_mapping_roundtrip_preserves_schema(store):
    from engine.pptx_mapper import load_template_mapping, save_template_mapping

    save_template_mapping("June Report.pptx", json.loads(json.dumps(SAMPLE_MAPPING)))
    loaded = load_template_mapping("June Report.pptx")

    assert loaded["_template_filename"] == "June Report.pptx"
    loaded.pop("_template_filename")
    assert loaded == SAMPLE_MAPPING


def test_mapping_filenames_are_stable_and_prefixed(store):
    """The mapping file name is derived from the template name; if this
    changes, every existing team mapping silently orphans."""
    from engine.pptx_mapper import save_template_mapping

    save_template_mapping("June Report.pptx", dict(SAMPLE_MAPPING))
    files = os.listdir(paths.MAPPINGS_DIR)

    assert len(files) == 1
    assert files[0].startswith("pptx_") and files[0].endswith(".json")


def test_resaving_updates_in_place_not_duplicates(store):
    from engine.pptx_mapper import load_template_mapping, save_template_mapping

    save_template_mapping("R.pptx", {"slides": {"1": {"shape_mappings": {}}}})
    save_template_mapping("R.pptx", dict(SAMPLE_MAPPING))

    assert len(os.listdir(paths.MAPPINGS_DIR)) == 1
    loaded = load_template_mapping("R.pptx")
    assert "1" in loaded["slides"]
    assert loaded["slides"]["1"]["shape_mappings"]["3"] == {"skip": True}


def test_missing_and_corrupt_mappings_return_none(store):
    from engine.pptx_mapper import (_mapping_path, load_template_mapping,
                                    save_template_mapping)

    assert load_template_mapping("never_saved.pptx") is None

    save_template_mapping("Broken.pptx", dict(SAMPLE_MAPPING))
    with open(_mapping_path("Broken.pptx"), "w", encoding="utf-8") as f:
        f.write("{ this is not json")
    assert load_template_mapping("Broken.pptx") is None  # no raise


def test_list_available_templates_reports_mapping_status(store):
    from engine.pptx_mapper import (list_available_templates,
                                    save_template_mapping)

    open(os.path.join(paths.TEMPLATES_DIR, "Mapped.pptx"), "wb").write(b"x")
    open(os.path.join(paths.TEMPLATES_DIR, "Unmapped.pptx"), "wb").write(b"x")
    open(os.path.join(paths.TEMPLATES_DIR, "notes.txt"), "w").write("skip me")
    save_template_mapping("Mapped.pptx", dict(SAMPLE_MAPPING))

    listed = {t["filename"]: t for t in list_available_templates()}

    assert set(listed) == {"Mapped.pptx", "Unmapped.pptx"}
    assert listed["Mapped.pptx"]["has_mapping"] is True
    assert listed["Unmapped.pptx"]["has_mapping"] is False
    assert listed["Mapped.pptx"]["path"] == os.path.join(
        paths.TEMPLATES_DIR, "Mapped.pptx")


# ── Template scanning ─────────────────────────────────────────────────────────

@pytest.fixture
def scanned_deck(tmp_path):
    """A one-slide deck with a title, plain textbox, table, and picture."""
    pytest.importorskip("PIL")
    from PIL import Image

    png = tmp_path / "pic.png"
    Image.new("RGB", (8, 8), (200, 0, 0)).save(png)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = "Monthly Recap"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    box.text_frame.text = "Impressions: [IMP]"
    slide.shapes.add_table(2, 3, Inches(1), Inches(3.5), Inches(4), Inches(1))
    slide.shapes.add_picture(str(png), Inches(6), Inches(2))
    path = str(tmp_path / "scan_me.pptx")
    prs.save(path)
    return path


def test_scan_reports_shape_types_and_positional_ids(scanned_deck):
    from engine.pptx_mapper import scan_template

    slides = scan_template(scanned_deck)

    assert len(slides) == 1
    slide = slides[0]
    assert slide["slide_num"] == 1
    assert slide["slide_title"] == "Monthly Recap"

    by_type = {}
    for shape in slide["shapes"]:
        by_type.setdefault(shape["shape_type"], []).append(shape)

    assert "Impressions: [IMP]" in [s["text"] for s in by_type["text"]]
    (table,) = by_type["table"]
    assert table["text"] == "[Table: 2r x 3c]"
    (image,) = by_type["image"]
    assert image["is_picture"] is True

    # shape_id remains the positional index (the mapping-JSON key and what
    # legacy mappings resolve by); shape_uid is the persistent PowerPoint
    # id that fills now prefer (engine.shape_identity).
    ids = [s["shape_id"] for s in slide["shapes"]]
    assert ids == sorted(ids)
    assert len(ids) == len(set(ids))
    uids = [s["shape_uid"] for s in slide["shapes"]]
    assert all(isinstance(u, int) for u in uids)
    assert len(uids) == len(set(uids))


def _shape_element(pptx_path, marker):
    """(presentation, lxml element) of the first shape whose text contains
    marker — element moves/removals preserve shape ids, exactly like a
    template edit in PowerPoint."""
    prs = Presentation(pptx_path)
    sp = next(s._element for s in prs.slides[0].shapes
              if s.has_text_frame and marker in s.text_frame.text)
    return prs, sp


def _mapping_with_identity(target, assignments):
    return {"slides": {"1": {"shape_mappings": {str(target["shape_id"]): {
        "shape_uid": target["shape_uid"], "shape_name": target["name"],
        "assignments": assignments}}}}}


def test_fill_follows_shape_identity_through_reordering(tmp_path, scanned_deck):
    """End-to-end drift test: map a shape (with its scanned uid), reorder
    the template so the positional index goes stale, and confirm the value
    still lands in the SAME shape. (Deliberately replaces the old
    positional-identity contract — mapper roadmap Phase 4.)"""
    from engine.pptx_fill import fill_template
    from engine.pptx_mapper import scan_template

    slides = scan_template(scanned_deck)
    target = next(s for s in slides[0]["shapes"]
                  if s["text"] == "Impressions: [IMP]")
    old_index = target["shape_id"]

    prs, sp = _shape_element(scanned_deck, "[IMP]")
    tree = sp.getparent()
    tree.remove(sp)
    tree.append(sp)  # move mapped shape to the end of the collection
    prs.save(scanned_deck)

    moved = next(s for s in scan_template(scanned_deck)[0]["shapes"]
                 if s["text"] == "Impressions: [IMP]")
    assert moved["shape_id"] != old_index      # index really drifted
    assert moved["shape_uid"] == target["shape_uid"]  # uid survived the edit

    out = str(tmp_path / "out.pptx")
    fill_template(scanned_deck, out,
                  _mapping_with_identity(target,
                      [{"metric": "Imps", "format": "text",
                        "replace_text": "[IMP]"}]),
                  {"Imps": "1,000,000"})

    texts = [s.text_frame.text for s in Presentation(out).slides[0].shapes
             if s.has_text_frame]
    assert "Impressions: 1,000,000" in texts
    assert not any("[IMP]" in t for t in texts)


def test_fill_survives_inserted_shape(tmp_path, scanned_deck):
    """A shape inserted before the mapped one shifts every positional index;
    the stored uid keeps the value on the right shape and the new shape
    untouched."""
    from engine.pptx_fill import fill_template_report
    from engine.pptx_mapper import scan_template

    target = next(s for s in scan_template(scanned_deck)[0]["shapes"]
                  if s["text"] == "Impressions: [IMP]")

    prs = Presentation(scanned_deck)
    box = prs.slides[0].shapes.add_textbox(Inches(1), Inches(0.2),
                                           Inches(2), Inches(0.5))
    box.text_frame.text = "NEW BANNER"
    sp = box._element
    tree = sp.getparent()
    tree.remove(sp)
    tree.insert(2, sp)  # spTree slots 0-1 are group properties, not shapes
    prs.save(scanned_deck)

    out = str(tmp_path / "out.pptx")
    _, report = fill_template_report(
        scanned_deck, out,
        _mapping_with_identity(target,
            [{"metric": "Imps", "format": "text", "replace_text": "[IMP]"}]),
        {"Imps": "42"})

    assert report.ok is True
    texts = [s.text_frame.text for s in Presentation(out).slides[0].shapes
             if s.has_text_frame]
    assert "Impressions: 42" in texts
    assert "NEW BANNER" in texts


def test_deleted_mapped_shape_is_reported_not_misfilled(tmp_path, scanned_deck):
    """Deleting a mapped shape must skip + report — never write into the
    shape that inherited its positional index."""
    from engine.pptx_fill import fill_template_report
    from engine.pptx_mapper import scan_template

    slides = scan_template(scanned_deck)
    title = next(s for s in slides[0]["shapes"]
                 if s["text"] == "Monthly Recap")
    target = next(s for s in slides[0]["shapes"]
                  if s["text"] == "Impressions: [IMP]")

    prs, sp = _shape_element(scanned_deck, "Monthly Recap")
    sp.getparent().remove(sp)  # target inherits the title's old index
    prs.save(scanned_deck)

    mapping = _mapping_with_identity(target,
        [{"metric": "Imps", "format": "text", "replace_text": "[IMP]"}])
    mapping["slides"]["1"]["shape_mappings"][str(title["shape_id"])] = {
        "shape_uid": title["shape_uid"], "shape_name": title["name"],
        "assignments": [{"metric": "T", "format": "text"}]}

    out = str(tmp_path / "out.pptx")
    _, report = fill_template_report(scanned_deck, out, mapping,
                                     {"T": "SHOULD NOT APPEAR", "Imps": "42"})

    assert report.ok is False
    assert report.missing_shapes == [title["name"]]
    texts = [s.text_frame.text for s in Presentation(out).slides[0].shapes
             if s.has_text_frame]
    assert "Impressions: 42" in texts        # sibling assignment still fills
    assert "SHOULD NOT APPEAR" not in texts  # nothing misfired


def test_scan_unreadable_file_raises_actionable_error(tmp_path):
    from engine.pptx_mapper import scan_template

    bad = tmp_path / "corrupt.pptx"
    bad.write_bytes(b"this is not a pptx at all")

    with pytest.raises(RuntimeError) as excinfo:
        scan_template(str(bad))
    assert "corrupt.pptx" in str(excinfo.value)
