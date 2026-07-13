"""Static slide-1 previews for saved PowerPoint templates.

A real PNG is exported through PowerPoint COM on Windows and cached under
``workspace/templates/thumbs/``.  When PowerPoint is unavailable, callers get
an immediate text summary from slide 1 instead of a permanently blank preview.
"""

from __future__ import annotations

import logging
import os

from config.paths import TEMPLATES_DIR

logger = logging.getLogger(__name__)

THUMBS_DIR = os.path.join(TEMPLATES_DIR, "thumbs")


def _thumb_path(template_path):
    base = os.path.splitext(os.path.basename(template_path))[0]
    return os.path.join(THUMBS_DIR, f"{base}.png")


def get_template_thumbnail(template_path, width=480):
    """Return a cached/generated slide-1 PNG path, or ``None``."""
    if not os.path.exists(template_path):
        return None

    png = _thumb_path(template_path)
    try:
        if (os.path.exists(png)
                and os.path.getmtime(png) >= os.path.getmtime(template_path)):
            return png
    except OSError:
        logger.debug("Thumb mtime check failed for %s", png, exc_info=True)

    return _export_thumbnail(template_path, png, width)


def get_template_preview(template_path, width=480):
    """Return preview data safe to compute on a background worker.

    The result is ``{"kind": "image", "value": <png path>}`` when COM export
    succeeds, otherwise ``{"kind": "text", "value": <slide summary>}``.
    Keeping PIL/Tk image creation out of this function lets the UI remain
    responsive while PowerPoint starts.
    """
    png = get_template_thumbnail(template_path, width=width)
    if png:
        return {"kind": "image", "value": png}
    return {"kind": "text", "value": _slide_text_summary(template_path)}


def _slide_text_summary(template_path, max_items=6, max_chars=42):
    """Return a concise slide-1 text fallback for preview panels."""
    filename = os.path.basename(template_path)
    try:
        from pptx import Presentation

        prs = Presentation(template_path)
        texts = []
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text[:max_chars])
        return "\n".join(texts[:max_items]) if texts else filename
    except Exception:
        logger.debug("Text preview failed for %s", template_path, exc_info=True)
        return filename


def _export_thumbnail(template_path, png_path, width):
    """Export slide 1 as PNG via an isolated PowerPoint COM instance."""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        logger.debug("pywin32 unavailable — no thumbnail for %s", template_path)
        return None

    os.makedirs(THUMBS_DIR, exist_ok=True)
    app = None
    prs = None
    initialized = False
    try:
        pythoncom.CoInitialize()
        initialized = True
        # DispatchEx avoids attaching to and later closing a PowerPoint window
        # the user already had open.
        app = win32com.client.DispatchEx("PowerPoint.Application")
        prs = app.Presentations.Open(
            os.path.abspath(template_path), ReadOnly=True, WithWindow=False)
        if prs.Slides.Count < 1:
            return None
        prs.Slides(1).Export(os.path.abspath(png_path), "PNG", width)
        return png_path if os.path.exists(png_path) else None
    except Exception:
        logger.debug("Thumbnail export failed for %s", template_path,
                     exc_info=True)
        return None
    finally:
        if prs is not None:
            try:
                prs.Close()
            except Exception:
                logger.debug("Thumbnail presentation close failed", exc_info=True)
        if app is not None:
            try:
                app.Quit()
            except Exception:
                logger.debug("Thumbnail PowerPoint close failed", exc_info=True)
        if initialized:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass
