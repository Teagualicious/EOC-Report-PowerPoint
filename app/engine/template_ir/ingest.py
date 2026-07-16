"""Phase 1 of the template-first pipeline: PPTX → template IR + assets.

Every shape's ORIGINAL XML is stored verbatim (the build-time render
source); images referenced by the XML are extracted into the template's
assets/ directory and recorded per relationship id so the builder can
re-target them. Chart shapes additionally get their chart PART extracted
(chart XML + embedded workbook + colors/style parts) so the builder can
clone the part wholesale — per the review doc, cloning is the only path
that keeps chart styling intact. A chart whose part cannot be extracted
is marked unsupported and skipped loudly at build.

Template stores live under workspace/ (client branding never enters the
repo — root CLAUDE.md data hygiene).
"""

import logging
import os
import re
from datetime import datetime, timezone

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from config.naming import storage_key
from config.paths import WORKSPACE_DIR
from engine.template_ir.schema import (ASSETS_DIR_NAME, SCHEMA_VERSION,
                                       ShapeIR, SlideIR, TemplateIR,
                                       save_template_ir)

logger = logging.getLogger(__name__)

TEMPLATE_STORE_DIR = os.path.join(WORKSPACE_DIR, "template_store")
SOURCE_PPTX_NAME = "source.pptx"   # kept for slide previews / re-checks

_CHART_TAG = qn("c:chart")
_TABLE_TAG = qn("a:tbl")
_BLIP_TAG = qn("a:blip")
_EMBED_ATTR = qn("r:embed")
_LINK_ATTR = qn("r:link")
_RID_ATTR = qn("r:id")


def ingest_template(pptx_path, template_id=None, store_dir=None):
    """Ingest a deck into the template store. Returns the template dir.

    store_dir overrides the default workspace/template_store/<template_id>/
    (tests and future UI pass it explicitly)."""
    template_id = template_id or storage_key(
        os.path.splitext(os.path.basename(pptx_path))[0], replace_dots=True)
    template_dir = store_dir or os.path.join(TEMPLATE_STORE_DIR, template_id)
    assets_dir = os.path.join(template_dir, ASSETS_DIR_NAME)
    os.makedirs(assets_dir, exist_ok=True)

    # Keep the source deck in the store: slide previews render from it,
    # and it documents exactly what was ingested.
    source_copy = os.path.join(template_dir, SOURCE_PPTX_NAME)
    if os.path.abspath(pptx_path) != os.path.abspath(source_copy):
        try:
            import shutil
            shutil.copy2(pptx_path, source_copy)
        except OSError:
            logger.warning("Could not keep a source copy in %s",
                           template_dir, exc_info=True)

    prs = Presentation(pptx_path)
    slides = []
    for slide_idx, slide in enumerate(prs.slides):
        shapes = []
        for z, shape in enumerate(slide.shapes):
            shapes.append(_ingest_shape(shape, slide, slide_idx, z, assets_dir))
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            layout_name = ""
        slides.append(SlideIR(slide_index=slide_idx, layout_name=layout_name,
                              shapes=shapes))

    ir = TemplateIR(
        schema_version=SCHEMA_VERSION,
        template_id=template_id,
        template_name=os.path.splitext(os.path.basename(pptx_path))[0],
        source_file=os.path.basename(pptx_path),
        created_at=datetime.now(timezone.utc).isoformat(timespec="seconds"),
        slide_width_emu=int(prs.slide_width),
        slide_height_emu=int(prs.slide_height),
        slides=slides,
    )
    save_template_ir(ir, template_dir)
    logger.info("Template ingested: %s -> %s (%d slide(s))",
                pptx_path, template_dir, len(slides))
    return template_dir


def _ingest_shape(shape, slide, slide_idx, z_order, assets_dir):
    element = shape._element
    xml = etree.tostring(element).decode("utf-8")

    try:
        uid = shape.shape_id
    except Exception:
        uid = None
    shape_id = f"s{slide_idx}_u{uid if uid is not None else f'z{z_order}'}"

    geometry = {}
    for prop in ("left", "top", "width", "height"):
        try:
            val = getattr(shape, prop)
            geometry[prop] = int(val) if val is not None else None
        except Exception:
            geometry[prop] = None
    try:
        geometry["rotation"] = float(shape.rotation)
    except Exception:
        geometry["rotation"] = 0.0

    text = ""
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text
    except Exception:
        logger.debug("Text snapshot failed for %s", shape_id, exc_info=True)

    shape_type, unsupported = _classify_element(element, shape)

    chart_part = {}
    if shape_type == "chart":
        chart_part = _extract_chart_part(element, slide, assets_dir, shape_id)
        if not chart_part:
            unsupported = "chart part could not be extracted"

    # Extract every embedded image the XML references (covers pictures,
    # picture fills, and pictures nested inside groups) and remember which
    # relationship id used it so the builder can re-target.
    image_rels = {}
    for blip in element.iter(_BLIP_TAG):
        rid = blip.get(_EMBED_ATTR)
        if not rid:
            if blip.get(_LINK_ATTR):
                unsupported = unsupported or "linked (external) picture"
            continue
        if rid in image_rels:
            continue
        asset = _extract_image(slide, rid, assets_dir, shape_id, len(image_rels))
        if asset:
            image_rels[rid] = asset
        else:
            unsupported = unsupported or f"unresolvable image relationship {rid}"

    return ShapeIR(shape_id=shape_id, shape_uid=uid, name=shape.name,
                   shape_type=shape_type, z_order=z_order, geometry=geometry,
                   element_xml=xml, text=text, image_rels=image_rels,
                   chart_part=chart_part, unsupported=unsupported)


def _classify_element(element, shape):
    """(shape_type, unsupported_reason_or_None) for one shape element."""
    tag = etree.QName(element).localname
    if tag == "graphicFrame":
        if element.find(f".//{_CHART_TAG}") is not None:
            return "chart", None       # part extracted separately (Phase C)
        if element.find(f".//{_TABLE_TAG}") is not None:
            return "table", None       # a:tbl is self-contained — copies fine
        return "other", "graphicFrame content (SmartArt/OLE?)"
    if tag == "pic":
        return "image", None
    if tag == "grpSp":
        return "group", None
    try:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return "text_box", None
    except Exception:
        logger.debug("has_text_frame check failed", exc_info=True)
    return "shape", None


def _extract_image(slide, rid, assets_dir, shape_id, seq):
    """Save the image part behind ``rid`` into assets/; return its relative
    path (or None). Filenames are content-addressed-ish (shape + seq +
    original extension) so re-ingest overwrites deterministically."""
    try:
        part = slide.part.related_part(rid)
        ext = (part.partname.ext or "png").lstrip(".")
        filename = f"{shape_id}_{seq}.{ext}"
        with open(os.path.join(assets_dir, filename), "wb") as f:
            f.write(part.blob)
        return f"{ASSETS_DIR_NAME}/{filename}"
    except Exception:
        logger.warning("Could not extract image %s for %s", rid, shape_id,
                       exc_info=True)
        return None


def _extract_chart_part(element, slide, assets_dir, shape_id):
    """Save the chart part behind a chart graphicFrame — its XML, embedded
    workbook, and auxiliary parts (chartColors/chartStyle) — into assets/.
    Returns the ShapeIR.chart_part dict, or {} when extraction fails."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    try:
        cchart = element.find(f".//{_CHART_TAG}")
        rid = cchart.get(_RID_ATTR)
        part = slide.part.related_part(rid)

        def save(suffix, blob):
            filename = f"{shape_id}_{suffix}"
            with open(os.path.join(assets_dir, filename), "wb") as f:
                f.write(blob)
            return f"{ASSETS_DIR_NAME}/{filename}"

        info = {"xml": save("chart.xml", part.blob), "workbook": None,
                "extra": []}
        n = 0
        for rel in part.rels.values():
            if rel.is_external:
                continue
            if rel.reltype == RT.PACKAGE:
                info["workbook"] = save("chart.xlsx", rel.target_part.blob)
            else:                       # chartColors/chartStyle etc.
                n += 1
                info["extra"].append(
                    [rel.reltype, rel.target_part.content_type,
                     save(f"chart_aux{n}.xml", rel.target_part.blob)])
        return info
    except Exception:
        logger.warning("Could not extract chart part for %s", shape_id,
                       exc_info=True)
        return {}


def list_template_stores(store_root=None):
    """Absolute paths of every ingested template store (dirs holding a
    template.json), sorted by name. Used by the template-first GUIs."""
    from engine.template_ir.schema import TEMPLATE_JSON_NAME
    root = store_root or TEMPLATE_STORE_DIR
    if not os.path.isdir(root):
        return []
    return [os.path.join(root, name) for name in sorted(os.listdir(root))
            if os.path.isfile(os.path.join(root, name, TEMPLATE_JSON_NAME))]


def strip_embed_ids(xml):
    """Normalize relationship ids out of shape XML so two ingests of the
    same visual content compare equal (rIds are packaging details that
    legitimately change across rebuilds). Used by round-trip tests."""
    return re.sub(r'r:(embed|id)="[^"]*"', r'r:\1=""', xml)
