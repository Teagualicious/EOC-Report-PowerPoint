"""Phase 4 of the template-first pipeline: template IR → new PPTX.

The builder never re-authors shapes from parsed properties: every
supported shape is a VERBATIM copy of its original XML (pixel-exact by
construction), with image relationships re-targeted to the extracted
assets. Unsupported shapes are skipped and REPORTED — silent partial
output is the failure mode this project has been burned by before.

Phase B: dynamic text. After the verbatim copies land, slot values are
written into their placeholder runs via engine.pptx_fill's replace
machinery (same multi-line/<a:br/> handling and outcome-based reporting
the production fill path uses). Only run TEXT changes — formatting stays
verbatim by construction.

Phase C: charts and tables. Chart PARTS are cloned wholesale (chart XML
+ embedded workbook + colors/style parts, relationships rewired) and
data injected through python-pptx's ``chart.replace_data`` — which
replaces only each series' cached tx/cat/val, leaving all series-level
formatting untouched. Table slots fill cells in place (the template's
header row is kept), cloning the last data row on overflow and clearing
leftovers.
"""

import logging
import os
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.parts.chart import ChartPart, EmbeddedXlsxPart
from pptx.util import Emu

from engine.template_ir.schema import load_template_ir

logger = logging.getLogger(__name__)

_BLIP_TAG = qn("a:blip")
_EMBED_ATTR = qn("r:embed")
_CHART_TAG = qn("c:chart")
_RID_ATTR = qn("r:id")


def build_from_template(template_dir, output_path, slot_values=None):
    """Rebuild a deck from a template store directory.

    ``slot_values`` maps slot names (from the IR's slot registry) to
    display strings; each is written into its placeholder run after the
    verbatim copy. Returns (output_path, report) where report is a
    JSON-friendly dict: shapes_copied, images_relinked,
    skipped: [(shape_id, reason)], slots_filled,
    slots_unfilled: [(slot_name, reason)].
    """
    ir = load_template_ir(template_dir)
    report = {"template_id": ir.template_id, "shapes_copied": 0,
              "images_relinked": 0, "charts_cloned": 0, "skipped": [],
              "slots_filled": 0, "slots_unfilled": [], "notes": []}

    prs = Presentation()
    prs.slide_width = Emu(ir.slide_width_emu)
    prs.slide_height = Emu(ir.slide_height_emu)
    blank = prs.slide_layouts[6]

    for slide_ir in ir.slides:
        slide = prs.slides.add_slide(blank)
        sp_tree = slide.shapes._spTree
        appended = {}
        for shape_ir in sorted(slide_ir.shapes, key=lambda s: s.z_order):
            if shape_ir.excluded:
                report["skipped"].append((shape_ir.shape_id,
                                          "excluded in template review"))
                continue
            if shape_ir.unsupported:
                report["skipped"].append((shape_ir.shape_id,
                                          shape_ir.unsupported))
                logger.warning("Build skipped %s: %s", shape_ir.shape_id,
                               shape_ir.unsupported)
                continue
            # python-pptx's typed parser (not raw lxml) so the appended
            # element can be wrapped by slide.shapes for the slot fill
            element = parse_xml(shape_ir.element_xml)
            report["images_relinked"] += _relink_images(
                element, shape_ir, slide, template_dir, report)
            if shape_ir.chart_part:
                if not _clone_chart_part(element, shape_ir, slide,
                                         template_dir, report):
                    continue   # reported by _clone_chart_part
            sp_tree.append(element)
            appended[shape_ir.shape_id] = element
            report["shapes_copied"] += 1
        if slot_values:
            _fill_slide_slots(slide, slide_ir, ir.slot_registry,
                              slot_values, appended, report)

    if slot_values is not None:
        registry = ir.slot_registry
        for slot in slot_values:
            if slot not in registry:
                report["slots_unfilled"].append(
                    (slot, "not a slot in this template"))
        for slot in registry:
            if slot not in slot_values:
                report["slots_unfilled"].append((slot, "no value provided"))

    prs.save(output_path)
    logger.info("Template rebuilt: %s -> %s | copied=%d relinked=%d skipped=%d "
                "slots=%d/%d",
                ir.template_id, output_path, report["shapes_copied"],
                report["images_relinked"], len(report["skipped"]),
                report["slots_filled"],
                report["slots_filled"] + len(report["slots_unfilled"]))
    return output_path, report


def _fill_slide_slots(slide, slide_ir, registry, slot_values, appended,
                      report):
    """Write slot values into their placeholder runs on one built slide.

    Resolution is by the slot's stored placeholder text inside its own
    shape (identity + text, per the review doc) — the shape was located
    by shape_id, so identical placeholders in different shapes can never
    cross. Failed matches are reported, never silent.
    """
    from engine.pptx_fill import _replace_in_text_frame
    shapes_by_id = {s.shape_id: s for s in slide_ir.shapes}
    for slot, spec in registry.items():
        if spec.get("slide_index") != slide_ir.slide_index \
                or slot not in slot_values:
            continue
        shape_ir = shapes_by_id.get(spec.get("shape_id"))
        if shape_ir is None:
            report["slots_unfilled"].append((slot, "shape not in template"))
            continue
        element = appended.get(shape_ir.shape_id)
        if element is None:
            reason = ("shape excluded in template review" if shape_ir.excluded
                      else f"shape skipped: {shape_ir.unsupported}")
            report["slots_unfilled"].append((slot, reason))
            continue
        wrapper = _wrap_shape(slide, element)
        slot_type = spec.get("type", "text")
        if slot_type == "chart_data":
            outcome = _inject_chart_data(wrapper, slot_values[slot])
        elif slot_type == "table_data":
            outcome = _fill_table(wrapper, slot_values[slot], report, slot)
        else:
            display = str(slot_values[slot])
            placeholder = spec.get("placeholder_text") or None
            outcome = None
            for text_frame in _iter_text_frames(wrapper):
                if _replace_in_text_frame(text_frame, display, placeholder):
                    break
            else:
                outcome = "placeholder text not found in shape"
        if outcome is None:
            report["slots_filled"] += 1
        else:
            report["slots_unfilled"].append((slot, outcome))


def _wrap_shape(slide, element):
    """The python-pptx wrapper for an element just appended to the slide's
    shape tree (identity match — the wrappers are built over the same
    lxml elements)."""
    for shape in slide.shapes:
        if shape._element is element:
            return shape
    return None


def _iter_text_frames(shape):
    """Yield every text frame under a shape, descending into groups."""
    if shape is None:
        return
    try:
        is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except NotImplementedError:   # python-pptx: unrecognized shape kinds
        is_group = False
    if is_group:
        for member in shape.shapes:
            yield from _iter_text_frames(member)
        return
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame


def _clone_chart_part(element, shape_ir, slide, template_dir, report):
    """Clone the stored chart part into the new package — chart XML plus
    embedded workbook and auxiliary (chartColors/chartStyle) parts, all
    relationships rewired — and point the copied graphicFrame at it.
    Returns True on success; on failure the shape is reported skipped."""
    info = shape_ir.chart_part
    package = slide.part.package
    try:
        with open(os.path.join(template_dir, info["xml"]), "rb") as f:
            chart_blob = f.read()
        chart_part = ChartPart.load(
            package.next_partname(ChartPart.partname_template),
            CT.DML_CHART, package, chart_blob)

        external_data = chart_part._element.externalData
        if info.get("workbook"):
            with open(os.path.join(template_dir, info["workbook"]), "rb") as f:
                xlsx_blob = f.read()
            xlsx_part = EmbeddedXlsxPart.load(
                package.next_partname(EmbeddedXlsxPart.partname_template),
                CT.SML_SHEET, package, xlsx_blob)
            rid = chart_part.relate_to(xlsx_part, RT.PACKAGE)
            chart_part._element.get_or_add_externalData().rId = rid
        elif external_data is not None:
            # No stored workbook: a dangling externalData rId would make
            # PowerPoint "repair" the deck — drop the element instead
            external_data.getparent().remove(external_data)

        for reltype, content_type, asset_rel in info.get("extra", []):
            with open(os.path.join(template_dir, asset_rel), "rb") as f:
                blob = f.read()
            aux = Part(package.next_partname("/ppt/charts/chartAux%d.xml"),
                       content_type, package, blob)
            chart_part.relate_to(aux, reltype)

        new_rid = slide.part.relate_to(chart_part, RT.CHART)
        element.find(f".//{_CHART_TAG}").set(_RID_ATTR, new_rid)
        report["charts_cloned"] += 1
        return True
    except Exception:
        logger.exception("Chart part clone failed for %s", shape_ir.shape_id)
        report["skipped"].append(
            (shape_ir.shape_id, "chart part clone failed"))
        return False


def _inject_chart_data(wrapper, payload):
    """Replace a cloned chart's data with a chart payload
    ({"categories", "series"}) via chart.replace_data — cached values and
    the embedded workbook change, series formatting does not. Returns
    None on success, else the unfilled reason."""
    from pptx.chart.data import CategoryChartData
    if wrapper is None or not getattr(wrapper, "has_chart", False):
        return "shape is not a chart"
    chart = wrapper.chart
    try:
        type_name = chart.chart_type.name
    except Exception:
        type_name = "unrecognized chart type"
    if type_name.startswith(("XY_", "BUBBLE")) or "unrecognized" in type_name:
        return f"data injection not supported for {type_name}"
    categories = payload.get("categories") or []
    series = payload.get("series") or []
    if not categories or not series:
        return "chart payload has no categories/series"
    # Keep the workbook's number format from the original cached values so
    # source-linked data labels don't drift to "General"
    fmt = chart._chartSpace.findtext(
        f".//{qn('c:val')}//{qn('c:formatCode')}") or "General"
    chart_data = CategoryChartData(number_format=fmt)
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name", ""), s.get("values", []))
    try:
        chart.replace_data(chart_data)
    except Exception as e:
        logger.exception("Chart data injection failed")
        return f"chart data injection failed: {e}"
    return None


def _fill_table(wrapper, payload, report, slot):
    """Write a table payload ({"headers", "rows"}) into a copied table.

    The template's own (branded) header row is KEPT — data rows fill from
    row 1 down. Overflow clones the last row (formatting preserved);
    leftover data rows are cleared. Column overflow is truncated and
    noted in the report. Returns None on success, else the reason."""
    from engine.pptx_fill import _replace_in_text_frame
    if wrapper is None or not getattr(wrapper, "has_table", False):
        return "shape is not a table"
    rows = payload.get("rows") or []
    if not rows:
        return "table payload has no rows"
    table = wrapper.table
    tbl = table._tbl
    n_cols = len(table.columns)
    trs = tbl.findall(qn("a:tr"))
    needed = 1 + len(rows)              # header row + data rows
    while len(trs) < needed:
        clone = deepcopy(trs[-1])
        trs[-1].addnext(clone)
        trs.append(clone)

    dropped = max(len(r) for r in rows) - n_cols
    if dropped > 0:
        report["notes"].append(
            f"{slot}: table has {n_cols} column(s); "
            f"{dropped} data column(s) dropped")

    for ri in range(1, len(trs)):
        data = rows[ri - 1] if ri - 1 < len(rows) else None
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            text = str(data[ci]) if data and ci < len(data) else ""
            if cell.text_frame.text == text:
                continue
            if not _replace_in_text_frame(cell.text_frame, text, None) \
                    and text:
                cell.text_frame.text = text
    return None


def _relink_images(element, shape_ir, slide, template_dir, report):
    """Point every image reference in the copied XML at a freshly added
    image part on the new slide. Returns how many were re-targeted."""
    if not shape_ir.image_rels:
        return 0
    relinked = 0
    new_rids = {}
    for old_rid, asset_rel in shape_ir.image_rels.items():
        asset_path = os.path.join(template_dir, asset_rel)
        try:
            _image_part, rid = slide.part.get_or_add_image_part(asset_path)
            new_rids[old_rid] = rid
        except Exception:
            logger.warning("Could not add image asset %s for %s", asset_rel,
                           shape_ir.shape_id, exc_info=True)
            report["skipped"].append(
                (shape_ir.shape_id, f"image asset missing: {asset_rel}"))
    for blip in element.iter(_BLIP_TAG):
        old = blip.get(_EMBED_ATTR)
        if old and old in new_rids:
            blip.set(_EMBED_ATTR, new_rids[old])
            relinked += 1
    return relinked
