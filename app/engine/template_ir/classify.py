"""Phase 2 of the template-first pipeline: suggest static/dynamic per
shape and propose named slots for placeholder-looking text runs.

Everything here produces SUGGESTIONS — the template review GUI presents
them for human confirmation before a template is used (proposal rule:
auto-classification is never final). All pure Python + lxml on the stored
shape XML, so the whole phase is exercised headlessly by the suite.

Slots pin to (shape identity + the run's placeholder text) per the review
doc — paragraph/run indices are stored as reconciliation hints only. Chart
and table shapes classify as dynamic but get NO slots until Phase C.
"""

import logging
import re

from lxml import etree
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

_P_TAG = qn("a:p")
_R_TAG = qn("a:r")
_T_TAG = qn("a:t")

# (compiled pattern, reason, slot type) — matched against one run's text.
_PLACEHOLDER_PATTERNS = [
    (re.compile(r"\{\{[^}]+\}\}"), "token placeholder", "text"),
    (re.compile(r"^\[[^\]]+\]$"), "bracketed placeholder", "text"),
    # Metric masks: X,XXX / XX% / $X,XXX / X.X — at least two X's or an
    # X with a numeric separator, so stray single letters never match.
    (re.compile(r"^\$?[Xx]{1,3}(?:[,.][Xx]{1,3})+%?$"),
     "metric mask (X,XXX)", "number"),
    (re.compile(r"^\$?[Xx]{2,}%?$"), "metric mask (XX)", "number"),
    (re.compile(r"^N/?A$", re.IGNORECASE), "N/A placeholder", "text"),
    (re.compile(r"\bCLIENT\s*NAME\b", re.IGNORECASE),
     "client-name field", "text"),
    (re.compile(r"\bMONTH\b.*\b(?:\d{4}|YYYY)\b", re.IGNORECASE),
     "date placeholder", "date"),
]

_MONTHS = (r"Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|"
           r"Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|"
           r"Nov(?:ember)?|Dec(?:ember)?")
_DATE_TEXT = re.compile(
    rf"\b(?:{_MONTHS})\b[\s.]*\d{{1,2}}(?:st|nd|rd|th)?,?\s+\d{{4}}",
    re.IGNORECASE)


def looks_like_placeholder(text):
    """(reason, slot_type) when the text looks like a data placeholder,
    else (None, None)."""
    t = (text or "").strip()
    if not t:
        return None, None
    for pattern, reason, slot_type in _PLACEHOLDER_PATTERNS:
        if pattern.search(t):
            return reason, slot_type
    if _DATE_TEXT.search(t):
        return "date text (variable header field)", "date"
    return None, None


def extract_runs(element_xml):
    """Flatten a shape's text runs from its stored XML, in document order
    (descends into groups). Returns [{"paragraph", "run", "text"}, ...]."""
    element = etree.fromstring(element_xml.encode("utf-8"))
    runs = []
    for p_idx, para in enumerate(element.iter(_P_TAG)):
        for r_idx, run in enumerate(para.findall(_R_TAG)):
            node = run.find(_T_TAG)
            runs.append({"paragraph": p_idx, "run": r_idx,
                         "text": node.text or "" if node is not None else ""})
    return runs


def _is_label_run(text):
    """All-caps alphabetic text like 'IMPRESSIONS' — the static half of a
    KPI callout pair, and the best source for the value's slot name."""
    t = (text or "").strip()
    return (len(t) >= 3 and t.isupper()
            and re.fullmatch(r"[A-Z][A-Z /&-]*[A-Z]", t) is not None
            and looks_like_placeholder(t) == (None, None))


def sanitize_slot_name(text, max_words=4):
    """'IMPRESSIONS' -> 'impressions'; 'Zip Code Chart!' -> 'zip_code_chart'."""
    words = re.findall(r"[A-Za-z0-9]+", (text or "").lower())[:max_words]
    name = "_".join(words)
    if name and name[0].isdigit():
        name = "slot_" + name
    return name


def _unique_name(base, used):
    name = base or "slot"
    n = 2
    while name in used:
        name = f"{base or 'slot'}_{n}"
        n += 1
    used.add(name)
    return name


def _suggest_slot_name(run, all_runs, reason):
    """Prefer the paired all-caps label run in the same shape (the run
    right after the placeholder, else right before); fall back to the
    reason/text itself."""
    idx = all_runs.index(run)
    for neighbor in (all_runs[idx + 1:idx + 2] + all_runs[max(0, idx - 1):idx]):
        if _is_label_run(neighbor["text"]):
            return sanitize_slot_name(neighbor["text"])
    if reason == "client-name field":
        return "client_name"
    if reason in ("date placeholder", "date text (variable header field)"):
        return "report_date"
    return sanitize_slot_name(run["text"])


def _classify_shape(shape):
    """(classification, reason, dynamic_runs, all_runs) for one ShapeIR."""
    if shape.shape_type == "chart":
        return "dynamic", "charts always receive mapped data", [], []
    if shape.shape_type == "table":
        return "dynamic", "tables always receive mapped data", [], []
    if shape.shape_type == "image":
        return "static", "logo or brand image", [], []

    runs = extract_runs(shape.element_xml)
    dynamic_runs = []
    reasons = []
    for run in runs:
        reason, slot_type = looks_like_placeholder(run["text"])
        if reason:
            dynamic_runs.append((run, reason, slot_type))
            if reason not in reasons:
                reasons.append(reason)
    if dynamic_runs:
        return "dynamic", "; ".join(reasons), dynamic_runs, runs
    if any(r["text"].strip() for r in runs):
        return "static", "no placeholder-like text", [], runs
    return "static", "decorative shape with no text", [], runs


def classify_template(ir):
    """Suggest classifications for every shape and rebuild the slot
    registry from placeholder-looking runs. Overwrites prior suggestions
    AND user edits — callers re-run it only on explicit request
    (re-ingest reconciliation that preserves review work is Phase D).
    Returns the ir for chaining."""
    ir.slot_registry = {}
    used = set()
    for slide in ir.slides:
        for shape in sorted(slide.shapes, key=lambda s: s.z_order):
            classification, reason, dynamic_runs, all_runs = \
                _classify_shape(shape)
            shape.classification = classification
            shape.classify_reason = reason
            shape.slot_name = None
            if shape.shape_type in ("chart", "table"):
                _add_frame_slot(ir, slide, shape, used)
                continue
            for run, run_reason, slot_type in dynamic_runs:
                name = _unique_name(
                    _suggest_slot_name(run, all_runs, run_reason), used)
                ir.slot_registry[name] = {
                    "type": slot_type,
                    "description": f"{run_reason}: \"{run['text'].strip()}\"",
                    "slide_index": slide.slide_index,
                    "shape_id": shape.shape_id,
                    "shape_uid": shape.shape_uid,
                    "placeholder_text": run["text"],
                    "paragraph": run["paragraph"],
                    "run": run["run"],
                }
                if shape.slot_name is None:
                    shape.slot_name = name
    return ir


def _add_frame_slot(ir, slide, shape, used):
    """One whole-shape slot for a chart or table (Phase C). Charts whose
    part could not be extracted get none — they cannot be built."""
    if shape.unsupported:
        return
    slot_type = "chart_data" if shape.shape_type == "chart" else "table_data"
    name = _unique_name(sanitize_slot_name(shape.name), used)
    ir.slot_registry[name] = {
        "type": slot_type,
        "description": f"{shape.shape_type} \"{shape.name}\" — map an "
                       "Advanced Query Builder query (Apply as "
                       f"{'Chart Data' if slot_type == 'chart_data' else 'Table'})",
        "slide_index": slide.slide_index,
        "shape_id": shape.shape_id,
        "shape_uid": shape.shape_uid,
        "placeholder_text": "",
        "paragraph": None,
        "run": None,
    }
    shape.slot_name = name


# ── Review-GUI mutations (pure, headlessly tested) ───────────────────────────

def shape_slots(ir, shape_id):
    """Slot names registered on one shape, in registry order."""
    return [name for name, spec in ir.slot_registry.items()
            if spec.get("shape_id") == shape_id]


def _find_shape(ir, shape_id):
    for slide in ir.slides:
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return slide, shape
    raise KeyError(f"unknown shape: {shape_id}")


def set_classification(ir, shape_id, classification):
    """User override from the review GUI. Switching a text shape to
    dynamic creates slots from its placeholder-looking runs (or its first
    non-empty run when none match, so there is always something to map);
    switching to static removes the shape's slots."""
    slide, shape = _find_shape(ir, shape_id)
    shape.classification = classification
    shape.classify_reason = "set in template review"
    if classification != "dynamic":
        for name in shape_slots(ir, shape_id):
            del ir.slot_registry[name]
        shape.slot_name = None
        return
    if shape_slots(ir, shape_id):
        return
    if shape.shape_type in ("chart", "table"):
        _add_frame_slot(ir, slide, shape, set(ir.slot_registry))
        return
    runs = extract_runs(shape.element_xml)
    used = set(ir.slot_registry)
    candidates = []
    for run in runs:
        reason, slot_type = looks_like_placeholder(run["text"])
        if reason:
            candidates.append((run, reason, slot_type))
    if not candidates:
        first = next((r for r in runs if r["text"].strip()), None)
        if first is None:
            return
        candidates = [(first, "marked dynamic in template review", "text")]
    for run, reason, slot_type in candidates:
        name = _unique_name(_suggest_slot_name(run, runs, reason), used)
        ir.slot_registry[name] = {
            "type": slot_type,
            "description": f"{reason}: \"{run['text'].strip()}\"",
            "slide_index": slide.slide_index,
            "shape_id": shape.shape_id,
            "shape_uid": shape.shape_uid,
            "placeholder_text": run["text"],
            "paragraph": run["paragraph"],
            "run": run["run"],
        }
        if shape.slot_name is None:
            shape.slot_name = name


def set_excluded(ir, shape_id, excluded):
    """Exclude a shape from the build entirely (its slots stay registered
    but report as unfilled — visible, never silent)."""
    _, shape = _find_shape(ir, shape_id)
    shape.excluded = bool(excluded)


def rename_slot(ir, old_name, new_name):
    """Rename a slot, preserving registry order. Raises ValueError on
    empty/duplicate/invalid names (GUI shows the message)."""
    new_name = sanitize_slot_name(new_name, max_words=8)
    if not new_name:
        raise ValueError("Slot name must contain letters or digits.")
    if new_name == old_name:
        return new_name
    if new_name in ir.slot_registry:
        raise ValueError(f"Slot name already in use: {new_name}")
    if old_name not in ir.slot_registry:
        raise KeyError(f"unknown slot: {old_name}")
    ir.slot_registry = {new_name if k == old_name else k: v
                        for k, v in ir.slot_registry.items()}
    for slide in ir.slides:
        for shape in slide.shapes:
            if shape.slot_name == old_name:
                shape.slot_name = new_name
    return new_name
