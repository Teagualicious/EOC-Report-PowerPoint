"""Fill a PowerPoint template with values, preserving all formatting."""

import logging
import os
import re
from copy import deepcopy

from pptx import Presentation
from pptx.oxml.ns import qn

from config.paths import APP_ROOT, WORKSPACE_DIR
from engine.pptx_formats import _coerce_number, _format_value, format_with_details
from engine.shape_identity import resolve_shape_index

logger = logging.getLogger(__name__)


def fill_template(template_path, output_path, mapping, metric_values):
    """Fill a PowerPoint template with actual values.

    Preserves ALL formatting: font, size, color, bold, italic, position, alignment.
    Only the text content changes — everything else stays identical.
    """
    path, _report = fill_template_report(template_path, output_path,
                                         mapping, metric_values)
    return path


def fill_template_report(template_path, output_path, mapping, metric_values):
    """Like fill_template, but also returns a FillReport describing what
    was filled, what was missing, and what failed. Returns (path, report).
    """
    from engine.fill_report import FillReport
    report = FillReport(template_path, output_path)
    prs = Presentation(template_path)

    for slide_num_str, slide_mapping in mapping.get("slides", {}).items():
        slide_idx = int(slide_num_str) - 1
        if slide_idx >= len(prs.slides):
            report.note_out_of_range()
            continue

        slide = prs.slides[slide_idx]
        shapes = list(slide.shapes)
        id_name_pairs = []
        for s in shapes:
            try:
                id_name_pairs.append((s.shape_id, s.name))
            except Exception:
                id_name_pairs.append((None, getattr(s, "name", "")))

        for shape_id_str, smap in slide_mapping.get("shape_mappings", {}).items():
            shape_idx = int(shape_id_str)

            if smap.get("skip", False):
                report.note_skip()
                continue

            # Prefer the stored persistent identity (uid, then unique name);
            # legacy entries without one resolve positionally as before.
            resolved = resolve_shape_index(id_name_pairs, shape_idx,
                                           smap.get("shape_uid"),
                                           smap.get("shape_name"))
            if resolved is None:
                if smap.get("shape_uid") is not None or smap.get("shape_name"):
                    # The mapped shape is gone from the template — skipping
                    # beats writing into whatever sits at the old index.
                    report.note_missing_shape(
                        smap.get("shape_name") or f"id {smap.get('shape_uid')}")
                else:
                    report.note_out_of_range()
                continue

            shape = shapes[resolved]

            # Image replacement — preserve position and size
            img_path = smap.get("image_path")
            if img_path:
                # Try absolute path first, then resolve relative from the app root
                resolved = img_path
                if not os.path.exists(resolved):
                    resolved = os.path.join(APP_ROOT, img_path)
                if not os.path.exists(resolved):
                    # Compatibility with mappings created before user data
                    # moved under workspace/ (e.g. templates/images/logo.png).
                    resolved = os.path.join(WORKSPACE_DIR, img_path)
                if not os.path.exists(resolved):
                    # Try the absolute path stored separately
                    resolved = smap.get("image_path_abs", "")
                if resolved and os.path.exists(resolved):
                    try:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        slide.shapes.add_picture(resolved, left, top, width, height)
                        sp = shape._element
                        sp.getparent().remove(sp)
                        report.note_image()
                    except Exception as img_exc:
                        logger.warning("Image replacement failed on slide %s shape %s",
                                       slide_num_str, shape_id_str, exc_info=True)
                        report.note_error(
                            f"image on slide {slide_num_str}: {img_exc}")
                else:
                    logger.warning("Image not found for slide %s shape %s: %s",
                                   slide_num_str, shape_id_str, img_path)
                    report.note_missing_image(img_path)
                continue

            # Build list of assignments (multi or single)
            assignments = smap.get("assignments", [])
            if not assignments and smap.get("metric"):
                assignments = [{"metric": smap["metric"],
                                 "format": smap.get("format", "text"),
                                 "replace_text": smap.get("replace_text"),
                                 "format_details": smap.get("format_details"),
                                 "query": smap.get("query")}]

            # Chart shapes: a saved Apply-as-Chart-Data query resolves to
            # categories/series and injects via python-pptx's cache
            # rewriter (series formatting untouched). Until now the
            # built-in engine had NO chart path — chart data only ever
            # reached the deck through the COM live preview, so Auto-Fill
            # silently left charts showing template numbers.
            if getattr(shape, "has_chart", False):
                for asgn in assignments:
                    query = asgn.get("query")
                    if not (isinstance(query, dict)
                            and query.get("output") == "chart"):
                        continue
                    from engine.query_resolver import resolve_query_payload
                    from engine.template_ir.build import inject_chart_data
                    payload = resolve_query_payload(
                        query, metric_values.get("__client_data__", []),
                        "chart")
                    if payload is None:
                        report.note_error(
                            f"chart on slide {slide_num_str}: query "
                            "matched no data rows")
                        continue
                    outcome = inject_chart_data(shape, payload)
                    if outcome is None:
                        report.note_filled()
                    else:
                        report.note_error(
                            f"chart on slide {slide_num_str}: {outcome}")
                continue

            for asgn in assignments:
                metric_name = asgn.get("metric", "")

                # Resolve query-based assignments
                query = asgn.get("query")
                if query and metric_name not in metric_values:
                    try:
                        from engine.query_resolver import resolve_query
                        val = resolve_query(query, metric_values.get("__client_data__", []),
                            metric_values.get("__client_name__", ""),
                            metric_values.get("__start_date__", ""),
                            metric_values.get("__end_date__", ""))
                        metric_values[metric_name] = val
                    except Exception:
                        logger.warning("Query resolution failed for %s", metric_name,
                                       exc_info=True)
                        report.note_failed_query(metric_name)

                if not metric_name or metric_name not in metric_values:
                    if metric_name and metric_name not in report.failed_queries:
                        report.note_missing_metric(metric_name)
                    continue

                value = metric_values[metric_name]
                fmt = asgn.get("format", "text")
                fmt_details = asgn.get("format_details")
                replace_text = asgn.get("replace_text")

                existing_text = ""
                if shape.has_text_frame:
                    existing_text = shape.text_frame.text
                # For partial replaces, detect date style / text case from the
                # TARGET substring, not the whole frame — otherwise unrelated
                # placeholders elsewhere in the line (e.g. "[MONTH/YEAR]")
                # hijack the format detection.
                fmt_context = replace_text if replace_text else existing_text

                # Detailed number formatting (decimals/commas/prefix/suffix)
                # wins for numeric values; standard formatter handles the
                # rest (dates, text case matching)
                coerced = _coerce_number(value)
                if fmt_details and fmt_details.get("format") == "date":
                    display = format_with_details(coerced, fmt_details)
                elif (fmt_details and isinstance(coerced, (int, float))
                        and not isinstance(coerced, bool)):
                    display = format_with_details(coerced, fmt_details)
                else:
                    display = _format_value(value, fmt, fmt_context)

                if shape.has_text_frame:
                    # Report from the replace's ACTUAL outcome. The old check
                    # matched replace_text against the frame's full text,
                    # which contains break characters (\n between paragraphs,
                    # \v for soft line breaks) that the runs being edited
                    # never do — multi-line targets were recorded as filled
                    # while the replace silently did nothing. This is the #1
                    # "my number didn't show up" failure mode.
                    if _replace_in_text_frame(shape.text_frame, display,
                                              replace_text):
                        report.note_filled()
                    else:
                        report.note_unmatched_placeholder(metric_name,
                                                          replace_text)

    prs.save(output_path)
    logger.info("PowerPoint filled: %s | %s | ok=%s", output_path,
                report.summary().splitlines()[0], report.ok)
    return output_path, report


def _replace_in_text_frame(text_frame, new_text, replace_text=None):
    """Replace text in a text frame while preserving all run-level formatting.

    If replace_text is specified, only that substring is replaced (partial).
    A multi-line target — the UI preview joins paragraphs with '\\n' and
    PowerPoint soft line breaks surface as '\\v', neither of which ever
    appears in the run text being searched — is matched line by line: the
    first target line receives the new text, the remaining target lines
    are cleared. If replace_text is None, the entire text content is
    replaced (an empty frame gets the text in its first paragraph).

    Returns True when something was written, so callers report unmatched
    placeholders from the actual outcome instead of guessing.
    """
    wrote = False
    if replace_text:
        lines = [ln.strip() for ln in re.split("[\r\n\v]+", replace_text)]
        lines = [ln for ln in lines if ln]
        if len(lines) <= 1:
            targets = [replace_text]
            if lines and lines[0] != replace_text:
                # Selection carried stray whitespace the runs don't have —
                # retry with the trimmed line before giving up
                targets.append(lines[0])
        else:
            targets = lines
        for target in targets:
            replacement = "" if wrote else new_text
            for para in text_frame.paragraphs:
                full_text = "".join(run.text for run in para.runs)
                if target not in full_text:
                    continue
                _replace_across_runs(para, target, replacement)
                if replacement and "\n" in replacement:
                    _explode_newlines(para)
                wrote = True
                break  # Only replace the first occurrence of each target
            if wrote and len(lines) <= 1:
                break  # single-line target placed — don't run the retry
    else:
        # Full replacement — replace text in the first paragraph that has content
        for para in text_frame.paragraphs:
            if para.text.strip():
                if para.runs:
                    # Put all new text in the first run, clear the rest
                    # This preserves the first run's formatting
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                    if "\n" in new_text:
                        _explode_newlines(para)
                else:
                    para.text = new_text
                wrote = True
                break
        if not wrote and text_frame.paragraphs:
            # An empty box assigned a full-text value used to no-op silently
            text_frame.paragraphs[0].text = new_text
            wrote = True
    return wrote


def _explode_newlines(paragraph):
    """Turn literal newlines inside runs into real <a:br/> line breaks.

    run.text stores '\\n' verbatim in <a:t>, which PowerPoint renders as
    whitespace rather than a break — a multi-line value looked right in
    the COM live preview but flattened (or vanished visually) in the
    exported deck. Each split keeps the source run's formatting.
    """
    for run in list(paragraph.runs):
        if "\n" not in run.text:
            continue
        parts = run.text.split("\n")
        r = run._r
        run.text = parts[0]
        anchor = r
        for part in parts[1:]:
            br = r.makeelement(qn("a:br"), {})
            anchor.addnext(br)
            new_r = deepcopy(r)
            new_r.find(qn("a:t")).text = part
            br.addnext(new_r)
            anchor = new_r


def _replace_across_runs(paragraph, old_text, new_text):
    """Replace old_text with new_text across paragraph runs, preserving formatting.

    Handles cases where the target text spans multiple runs.
    """
    runs = paragraph.runs
    if not runs:
        if old_text in paragraph.text:
            paragraph.text = paragraph.text.replace(old_text, new_text, 1)
        return

    # Build concatenated text and track run boundaries
    texts = [run.text for run in runs]
    full = "".join(texts)

    if old_text not in full:
        return

    idx = full.index(old_text)
    end_idx = idx + len(old_text)

    # Find which runs are affected
    pos = 0
    new_texts = []
    for i, run_text in enumerate(texts):
        run_start = pos
        run_end = pos + len(run_text)

        if run_end <= idx or run_start >= end_idx:
            # This run is entirely outside the replacement zone — keep as-is
            new_texts.append(run_text)
        elif run_start <= idx and run_end >= end_idx:
            # Entire replacement is within this single run
            before = run_text[:idx - run_start]
            after = run_text[end_idx - run_start:]
            new_texts.append(before + new_text + after)
        elif run_start <= idx < run_end:
            # Replacement starts in this run
            before = run_text[:idx - run_start]
            new_texts.append(before + new_text)
        elif run_start < end_idx <= run_end:
            # Replacement ends in this run
            after = run_text[end_idx - run_start:]
            new_texts.append(after)
        else:
            # This run is entirely within the replacement zone — clear it
            new_texts.append("")

        pos = run_end

    # Apply new texts to runs (preserving all formatting)
    for i, run in enumerate(runs):
        if i < len(new_texts):
            run.text = new_texts[i]
