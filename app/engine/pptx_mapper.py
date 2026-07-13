"""PowerPoint template scanning and mapping storage.

Companion modules (all re-exported here so callers can keep importing
from ``engine.pptx_mapper``):

* ``engine.pptx_formats`` — value/date/text formatting
* ``engine.pptx_fill`` — filling a template with values
* ``engine.metrics_catalog`` — building the available-metrics catalog
"""

import json
import logging
import os

from pptx import Presentation

from config.paths import MAPPINGS_DIR, TEMPLATES_DIR
from config.naming import storage_key
from config.settings import _atomic_json_write
# Re-exported for ui/mapper modules (public API of this module)
from engine.metrics_catalog import SPECIAL_METRICS, get_available_metrics  # noqa: F401
from engine.pptx_fill import fill_template  # noqa: F401

logger = logging.getLogger(__name__)


def scan_template(template_path):
    """Scan PowerPoint and return slide/shape structure."""
    try:
        return _scan_with_pptx(template_path)
    except Exception as e1:
        logger.warning("python-pptx scan failed for %s: %s", template_path, e1)
        # python-pptx failed — try with a temp copy
        try:
            import shutil, tempfile
            tmp = os.path.join(tempfile.gettempdir(), f"_scan_{os.path.basename(template_path)}")
            shutil.copy2(template_path, tmp)
            result = _scan_with_pptx(tmp)
            try:
                os.remove(tmp)
            except OSError:
                logger.debug("Could not remove temp scan copy %s", tmp)
            return result
        except Exception as e2:
            logger.warning("Temp-copy scan failed: %s", e2)
            # Last resort: try COM-based scan
            try:
                return _scan_with_com(template_path)
            except Exception as e3:
                logger.error("All scan methods failed for %s", template_path)
                raise RuntimeError(
                    f"Could not read the PowerPoint file "
                    f"'{os.path.basename(template_path)}'.\n"
                    f"pptx: {e1}\ncopy: {e2}\nCOM: {e3}") from e3


def _scan_with_pptx(template_path):
    """Scan using python-pptx."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(template_path)
    slides = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {"slide_num": slide_idx + 1, "slide_title": "", "shapes": []}
        if slide.shapes.title:
            slide_info["slide_title"] = slide.shapes.title.text

        for shape_idx, shape in enumerate(slide.shapes):
            is_picture = False
            try:
                is_picture = shape.shape_type in (MSO_SHAPE_TYPE.PICTURE,
                                                  MSO_SHAPE_TYPE.LINKED_PICTURE)
            except Exception:
                logger.debug("shape_type check failed for %s", shape.name)

            shape_info = {
                "name": shape.name, "shape_id": shape_idx,
                "has_text": shape.has_text_frame,
                "has_table": hasattr(shape, "has_table") and shape.has_table,
                "has_chart": hasattr(shape, "has_chart") and shape.has_chart,
                "is_picture": is_picture,
                "text": "", "full_text": "", "shape_type": "text",
            }

            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                shape_info["text"] = "\n".join(texts)
                shape_info["full_text"] = shape_info["text"]

            if hasattr(shape, "has_table") and shape.has_table:
                shape_info["shape_type"] = "table"
                shape_info["text"] = f"[Table: {len(shape.table.rows)}r x {len(shape.table.columns)}c]"

            if hasattr(shape, "has_chart") and shape.has_chart:
                shape_info["shape_type"] = "chart"
                shape_info["text"] = f"[Chart: {shape.name}]"

            if is_picture:
                shape_info["shape_type"] = "image"
                shape_info["text"] = f"[Image: {shape.name}]"

            if (shape_info["text"] or shape.has_text_frame
                    or shape_info["shape_type"] in ("table", "chart", "image")):
                slide_info["shapes"].append(shape_info)

        slides.append(slide_info)
    return slides


def _scan_with_com(template_path):
    """Fallback scan using COM automation (Windows only, requires pywin32)."""
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        raise RuntimeError("COM scan unavailable: pywin32 is not installed") from exc
    pythoncom.CoInitialize()

    app = win32com.client.Dispatch("PowerPoint.Application")
    prs = app.Presentations.Open(os.path.abspath(template_path), ReadOnly=True, WithWindow=False)

    slides = []
    for si in range(1, prs.Slides.Count + 1):
        slide = prs.Slides(si)
        slide_info = {"slide_num": si, "slide_title": "", "shapes": []}

        for shi in range(1, slide.Shapes.Count + 1):
            shape = slide.Shapes(shi)
            # msoPicture = 13, msoLinkedPicture = 11
            is_picture = False
            try:
                is_picture = shape.Type in (13, 11)
            except Exception:
                logger.debug("COM shape Type check failed")
            shape_info = {
                "name": shape.Name, "shape_id": shi - 1,
                "has_text": shape.HasTextFrame, "has_table": shape.HasTable if hasattr(shape, 'HasTable') else False,
                "has_chart": shape.HasChart if hasattr(shape, 'HasChart') else False,
                "is_picture": is_picture,
                "text": "", "full_text": "", "shape_type": "text",
            }

            if shape.HasTextFrame:
                text = shape.TextFrame.TextRange.Text
                shape_info["text"] = text.strip()
                shape_info["full_text"] = text
                if not slide_info["slide_title"] and shi <= 2:
                    slide_info["slide_title"] = text.strip()[:60]

            if hasattr(shape, 'HasTable') and shape.HasTable:
                shape_info["shape_type"] = "table"
                shape_info["text"] = f"[Table: {shape.Table.Rows.Count}r x {shape.Table.Columns.Count}c]"

            if hasattr(shape, 'HasChart') and shape.HasChart:
                shape_info["shape_type"] = "chart"
                shape_info["text"] = f"[Chart: {shape.Name}]"

            if is_picture:
                shape_info["shape_type"] = "image"
                shape_info["text"] = f"[Image: {shape.Name}]"

            if (shape_info["text"] or shape.HasTextFrame
                    or shape_info["shape_type"] in ("table", "chart", "image")):
                slide_info["shapes"].append(shape_info)

        slides.append(slide_info)

    prs.Close()
    pythoncom.CoUninitialize()
    return slides


def _mapping_path(template_filename):
    safe = storage_key(template_filename, replace_dots=True)
    return os.path.join(MAPPINGS_DIR, f"pptx_{safe}.json")


def load_template_mapping(template_filename):
    """Load the saved JSON mapping for a template, or None."""
    path = _mapping_path(template_filename)
    if os.path.exists(path):
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except (OSError, json.JSONDecodeError):
            logger.warning("Corrupt or unreadable mapping file: %s", path,
                           exc_info=True)
    return None


def save_template_mapping(template_filename, mapping):
    """Persist a template mapping as JSON."""
    os.makedirs(MAPPINGS_DIR, exist_ok=True)
    # Store template filename for reference
    mapping["_template_filename"] = template_filename
    _atomic_json_write(_mapping_path(template_filename), mapping)



def list_available_templates():
    """List all templates that have been imported, with mapping status."""
    if not os.path.exists(TEMPLATES_DIR):
        return []
    templates = []
    for f in os.listdir(TEMPLATES_DIR):
        if f.endswith(".pptx"):
            mapping = load_template_mapping(f)
            templates.append({
                "filename": f,
                "path": os.path.join(TEMPLATES_DIR, f),
                "has_mapping": mapping is not None,
            })
    return templates
